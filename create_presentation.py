"""
Dissertation Presentation Generator (PowerPoint)
Creates 30-minute presentation for committee review
Assumes committee has read proposal; focuses on visual findings
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_presentation():
    """Generate the dissertation presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define color scheme
    DARK_BLUE = RGBColor(31, 78, 121)
    LIGHT_BLUE = RGBColor(79, 129, 189)
    ACCENT_RED = RGBColor(192, 0, 0)
    ACCENT_GREEN = RGBColor(0, 176, 80)
    WHITE = RGBColor(255, 255, 255)
    DARK_GRAY = RGBColor(89, 89, 89)

    def add_title_slide(title, subtitle=""):
        """Add a title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BLUE

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            p = subtitle_frame.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = LIGHT_BLUE

        return slide

    def add_content_slide(title, bullet_points=None):
        """Add a content slide with title and bullets."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title box (dark blue background)
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = DARK_BLUE
        title_shape.line.color.rgb = DARK_BLUE

        title_frame = title_shape.text_frame
        title_frame.word_wrap = False
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Content box
        if bullet_points:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(6))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True

            for i, point in enumerate(bullet_points):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                p.text = point
                p.font.size = Pt(18)
                p.font.color.rgb = DARK_GRAY
                p.space_before = Pt(12)
                p.space_after = Pt(12)
                p.level = 0

        return slide

    def add_key_metric_slide(title, metric_title, metric_value, metric_unit="", description=""):
        """Add a slide highlighting a key metric."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title box
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = DARK_BLUE
        title_shape.line.color.rgb = DARK_BLUE

        title_frame = title_shape.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Large metric box
        metric_shape = slide.shapes.add_shape(1, Inches(1.5), Inches(1.5), Inches(7), Inches(2.5))
        metric_shape.fill.solid()
        metric_shape.fill.fore_color.rgb = LIGHT_BLUE
        metric_shape.line.color.rgb = DARK_BLUE
        metric_shape.line.width = Pt(3)

        # Metric label
        label_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.7), Inches(7), Inches(0.6))
        label_frame = label_box.text_frame
        p = label_frame.paragraphs[0]
        p.text = metric_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Metric value
        value_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.4), Inches(7), Inches(1.2))
        value_frame = value_box.text_frame
        p = value_frame.paragraphs[0]
        p.text = f"{metric_value} {metric_unit}"
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        if description:
            desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2.5))
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            p = desc_frame.paragraphs[0]
            p.text = description
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY

        return slide

    def add_comparison_slide(title, left_label, left_value, right_label, right_value, color_left=ACCENT_GREEN, color_right=ACCENT_RED):
        """Add a comparison slide with two metrics."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title box
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = DARK_BLUE
        title_shape.line.color.rgb = DARK_BLUE

        title_frame = title_shape.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Left box
        left_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(4.2), Inches(5))
        left_shape.fill.solid()
        left_shape.fill.fore_color.rgb = color_left
        left_shape.line.color.rgb = DARK_GRAY
        left_shape.line.width = Pt(2)

        left_label_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4.2), Inches(0.8))
        left_label_frame = left_label_box.text_frame
        left_label_frame.word_wrap = True
        p = left_label_frame.paragraphs[0]
        p.text = left_label
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE

        left_value_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(4.2), Inches(2.5))
        left_value_frame = left_value_box.text_frame
        left_value_frame.word_wrap = True
        p = left_value_frame.paragraphs[0]
        p.text = left_value
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Right box
        right_shape = slide.shapes.add_shape(1, Inches(5.3), Inches(1.5), Inches(4.2), Inches(5))
        right_shape.fill.solid()
        right_shape.fill.fore_color.rgb = color_right
        right_shape.line.color.rgb = DARK_GRAY
        right_shape.line.width = Pt(2)

        right_label_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.7), Inches(4.2), Inches(0.8))
        right_label_frame = right_label_box.text_frame
        right_label_frame.word_wrap = True
        p = right_label_frame.paragraphs[0]
        p.text = right_label
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE

        right_value_box = slide.shapes.add_textbox(Inches(5.3), Inches(3), Inches(4.2), Inches(2.5))
        right_value_frame = right_value_box.text_frame
        right_value_frame.word_wrap = True
        p = right_value_frame.paragraphs[0]
        p.text = right_value
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        return slide

    # SLIDE 1: Title
    add_title_slide(
        "Data Breach Disclosure Timing & Market Outcomes",
        "A Natural Experiment Analysis of FCC Cybersecurity Mandates"
    )

    # SLIDE 2: The Central Finding
    add_key_metric_slide(
        "The Core Discovery",
        "Markets Punish WHO YOU ARE & WHAT WAS BREACHED",
        "Not WHEN YOU TALK",
        "",
        "Disclosure timing has NO effect on shareholder returns (H1 null). \n"
        "Firm characteristics and breach severity drive market reactions. \n"
        "Mandatory timing mandates create costs without valuation benefits."
    )

    # SLIDE 3: Dataset
    add_content_slide(
        "Dataset: 1,054 Breaches (2006-2025)",
        [
            "926 breaches with CRSP stock price data (Essay 1: Market Returns)",
            "916 breaches with volatility data (Essay 2: Information Asymmetry)",
            "896 breaches with executive change data (Essay 3: Governance)",
            "200 FCC-regulated firms (19%); 854 non-FCC firms (81%)",
            "442 repeat offenders (41.9% with prior breach history)",
            "117 health breaches (11.1%), 257 financial (24.4%)",
            "Disclosure timing: 19% ≤7 days, 34% 8-30 days, 47% >30 days"
        ]
    )

    # SLIDE 4: Research Design
    add_content_slide(
        "Causal Identification: FCC Rule 37.3 Natural Experiment",
        [
            "FCC Rule effective January 1, 2007: 7-day disclosure mandate for SIC 4813/4899/4841",
            "Control group: all other industries under state-level laws (30-90 day timelines)",
            "Validation Test 1 (Temporal): FCC effect zero pre-2007 (p=0.88), significant post-2007 (p=0.0125) ✓",
            "Validation Test 2 (Industry Controls): Effect strengthens with controls (-2.20% → -5.37%) ✓",
            "Validation Test 3 (Size Sensitivity): Heterogeneous effects prove regulation, not size selection ✓",
            "27+ robustness specifications, machine learning validation, falsification tests all support findings"
        ]
    )

    # SLIDE 5: H1 - Timing Effect (Main Finding)
    add_key_metric_slide(
        "H1: Does Disclosure Timing Affect Shareholder Returns?",
        "Timing Coefficient",
        "+0.57%",
        "(p=0.539, not significant)",
        "ANSWER: NO. Disclosure timing is IRRELEVANT to market returns. \n"
        "Equivalence testing confirms this null is genuine (not power shortage). \n"
        "Refutes regulatory assumption that speed creates market benefits."
    )

    # SLIDE 6: H2 - FCC Regulatory Effect
    add_key_metric_slide(
        "H2: Does FCC Regulation Affect Returns?",
        "FCC Regulatory Penalty",
        "-2.20%",
        "CAR (p=0.010, significant)",
        "FCC-regulated firms experience ~$4.0M average shareholder loss per breach. \n"
        "Aggregate loss across 187 FCC breaches: $0.76 BILLION. \n"
        "Effect driven by regulatory burden signal, not firm selection (strengthens with controls)."
    )

    # SLIDE 7: H3 - Prior Breaches (Strongest Effect)
    add_key_metric_slide(
        "H3: Does Prior Breach History Matter?",
        "Reputational Accumulation (STRONGEST EFFECT)",
        "-0.22%",
        "per prior breach (p<0.001)",
        "A firm with 5 prior breaches faces -1.1% additional penalty on next breach. \n"
        "Markets view repeat breaches as evidence of SYSTEMATIC FAILURES, not isolated incidents. \n"
        "This is the dominant market discipline mechanism."
    )

    # SLIDE 8: H4 - Health Breaches
    add_key_metric_slide(
        "H4: Do Health Breaches Cost More?",
        "Health Breach Premium",
        "-2.51%",
        "CAR (p=0.004, significant)",
        "Health breaches produce nearly identical penalty magnitude to FCC regulatory penalty. \n"
        "Markets price in HIPAA complexity, FDA investigation risk, heightened liability. \n"
        "Regulatory complexity = market-equivalent cost to regulatory mandate."
    )

    # SLIDE 9: Comparison - What Drives Market Reactions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Essay 1: What Drives Market Reactions?"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(6))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    factors = [
        ("Prior Breaches", "-0.22% per breach***", "STRONGEST"),
        ("FCC Regulation", "-2.20%**", "Regulatory burden signal"),
        ("Health Breach", "-2.51%***", "Regulatory complexity"),
        ("Disclosure Timing", "+0.57%", "NOT SIGNIFICANT"),
        ("Breach Size", "Not significant", "Markets price fundamentals")
    ]

    for i, (factor, effect, interpretation) in enumerate(factors):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = f"{factor}: {effect} — {interpretation}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY if "NOT SIGNIFICANT" not in interpretation else ACCENT_RED
        p.space_before = Pt(10)
        p.space_after = Pt(10)

    # SLIDE 10: Essay 2 - Volatility Paradox
    add_key_metric_slide(
        "H5: Does FCC Regulation Increase Market Uncertainty?",
        "Volatility Increase",
        "+1.68% to +5.02%",
        "(Paradoxical!)",
        "PARADOX: Mandatory disclosure succeeds at SPEED but FAILS at QUALITY. \n"
        "Forced 7-day deadline prevents investigation completion. \n"
        "Incomplete disclosure (\"investigation ongoing\") INCREASES uncertainty instead of decreasing it. \n"
        "Size heterogeneity: Small firms +7.31% volatility; large firms -3.39% (capacity constraints)."
    )

    # SLIDE 11: Volatility Mechanism
    add_content_slide(
        "The Volatility Paradox: Why Mandatory Timing Increases Uncertainty",
        [
            "Regulatory Constraint: FCC 7-day deadline forces disclosure before investigation complete",
            "Information Quality Tradeoff: Complete investigation (45-60 days) vs. Speed (7 days)",
            "Market Recognition: Markets see incomplete disclosures as quality problems, raising uncertainty",
            "Size Matters: Small firms lack processing capacity to investigate rapidly + disclose completely",
            "Large firms can accommodate speed (more resources, better processes)",
            "Policy Implication: Speed mandate creates unintended information quality costs"
        ]
    )

    # SLIDE 12: Essay 3 - Governance Response
    add_key_metric_slide(
        "H6: Does Mandatory Disclosure Trigger Governance Changes?",
        "Executive Turnover Acceleration",
        "+5.3 percentage points",
        "(50.6% vs 45.3%)",
        "Immediate disclosure (≤7 days) accelerates executive departures. \n"
        "Mechanism: Stakeholder pressure → boards respond with governance changes. \n"
        "Turnover (46.4%) is 50x more common than regulatory enforcement (0.6%). \n"
        "This is organizational self-response mechanism, not formal regulation."
    )

    # SLIDE 13: Three Mechanisms Operating Independently
    add_content_slide(
        "Three Mechanisms Operate Independently",
        [
            "Essay 1 (Valuation): Timing IRRELEVANT. Markets price firm fundamentals, not announcement speed.",
            "Essay 2 (Uncertainty): Timing INCREASES volatility. Forced speed prevents investigation completion.",
            "Essay 3 (Governance): Timing ACCELERATES turnover. Stakeholder pressure drives response.",
            "",
            "MEDIATION TEST: Volatility does NOT mediate governance response (indirect effect p=0.48).",
            "Essays 2 and 3 operate through completely different channels.",
            "Regulation achieves governance activation at cost of increased uncertainty + reduced valuation."
        ]
    )

    # SLIDE 14: Economic Significance - FCC Cost
    add_comparison_slide(
        "Economic Significance: FCC Regulatory Costs",
        "Per Breach\n(Median Firm)",
        "-$0.9M",
        "Aggregate\n(187 breaches)",
        "-$0.76 BILLION",
        color_left=ACCENT_RED,
        color_right=ACCENT_RED
    )

    # SLIDE 15: Economic Significance - Volatility
    add_key_metric_slide(
        "Cost of Capital Impact: Information Asymmetry",
        "Volatility-Induced Cost of Capital Increase",
        "+$3-4M annually",
        "per $1B debt refinancing",
        "The +1.68% to +5.02% volatility increase translates to 39bp cost of capital increase. \n"
        "For firms refinancing large debt, this compounds to multi-million dollar annual costs. \n"
        "Effect is concentrated in small firms lacking investigation capacity."
    )

    # SLIDE 16: Economic Significance - Governance
    add_key_metric_slide(
        "Governance Disruption Costs",
        "Per Executive Departure",
        "$12-25M",
        "(direct + indirect costs)",
        "46.4% of breaches (416 incidents) experience 30-day executive turnover. \n"
        "Governance literature estimates total disruption cost: $2-5M direct + $10-20M indirect. \n"
        "Aggregate governance cost across sample: $0.39-0.98 BILLION."
    )

    # SLIDE 17: Total Economic Cost
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Total Economic Cost: FCC Regulation"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(6))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    costs = [
        "Valuation Loss (CAR effect): $0.76 Billion",
        "Cost of Capital Increases: $0.1-0.2 Billion (annualized over 10 years)",
        "Governance Disruption: $0.39-0.98 Billion",
        "Total Aggregate Cost: $1.25-1.94 BILLION"
    ]

    for i, cost in enumerate(costs):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = cost
        p.font.size = Pt(22)
        p.font.bold = True if "Total" in cost else False
        p.font.color.rgb = ACCENT_RED if "Total" in cost else DARK_GRAY
        p.space_before = Pt(12)
        p.space_after = Pt(12)

    # SLIDE 18: Policy Implications - FCC
    add_content_slide(
        "Policy Implications: FCC 7-Day Rule",
        [
            "Current Finding: Rule generates $0.76B shareholder losses without valuation benefits",
            "Root Cause: 7-day timeline forces disclosure before investigation completion",
            "Market Impact: Incomplete disclosure increases uncertainty instead of reducing it",
            "",
            "Recommendation 1: Extend to 14-21 days to allow basic investigation completion",
            "Recommendation 2: Allow conditional disclosure (\"preliminary findings pending investigation\")",
            "Recommendation 3: Focus on COMPLETENESS over SPEED for information quality"
        ]
    )

    # SLIDE 19: Policy Implications - SEC
    add_content_slide(
        "Policy Implications: SEC 4-Day Cybersecurity Rule (2023)",
        [
            "SEC explicitly requested empirical validation of 4-day timeline before implementation",
            "This dissertation provides that evidence: SEC faces identical risks to FCC",
            "4-day requirement likely to increase uncertainty while failing to change valuations",
            "",
            "Recommendation 1: Implement with explicit 18-24 month evaluation built in",
            "Recommendation 2: Budget IT systems for rapid investigation capability",
            "Recommendation 3: Monitor actual cost-benefit vs. regulatory intent"
        ]
    )

    # SLIDE 20: Policy Implications - Optimal Timeline
    add_content_slide(
        "The Optimal Disclosure Timeline Concept",
        [
            "Too Fast (7 days): Forces incomplete disclosure → increases uncertainty",
            "Too Slow (unregulated): Allows information asymmetry → information advantage",
            "Optimal (45-60 days): Allows investigation + keeps stakeholders informed",
            "",
            "Evidence Supporting 45-60 Day Window:",
            "• HIPAA uses 60 days (longer timeline than FCC's 7 days)",
            "• Stakeholder research (Xu et al. 2024): prefers completeness over speed",
            "• Our data: volatility effects concentrated in forced-speed scenarios, not moderate timelines"
        ]
    )

    # SLIDE 21: Firm Strategy
    add_content_slide(
        "Firm Disclosure Strategy Implications",
        [
            "Timing does NOT reduce shareholder losses (H1 null)",
            "Early disclosure IS valuable for governance credibility with stakeholders",
            "Optimal strategy: Disclose early for accountability, not for valuation protection",
            "",
            "Key insight: Early disclosure is a GOVERNANCE SIGNAL, not a MARKET SIGNAL",
            "Boards should use disclosure timing as accountability mechanism with stakeholders",
            "But do not expect disclosure speed to reduce shareholder market penalties"
        ]
    )

    # SLIDE 22: Validation Evidence
    add_content_slide(
        "Robustness & Validation: This Finding is Solid",
        [
            "✓ Mediation Analysis: Volatility does not mediate governance (independent mechanisms)",
            "✓ Event Window Sensitivity: Effects consistent across 5-day and 30-day windows",
            "✓ Falsification Tests: Effects are breach-specific, not general firm artifacts",
            "✓ Low R² Sensitivity: Alternative specs show no improvement (methodology adequate)",
            "✓ Machine Learning: Random Forest feature importance confirms prior_breaches #1, timing much lower",
            "✓ 27+ Robustness Specifications: Timing null holds across all variations"
        ]
    )

    # SLIDE 23: Limitations
    add_content_slide(
        "Limitations & Scope",
        [
            "Sample Composition: FCC firms 2x larger; controlled with size variables and subgroup analysis",
            "Scope: Evidence concerns shareholder returns only; other outcomes beyond scope",
            "Public Firms: Results generalize to publicly-traded firms with CRSP/EDGAR data",
            "Causal Chain: Parallel trends assumption validates temporally; pre-2007 test confirms",
            "Turnover Windows: 30/90/180-day analysis may miss longer-run effects",
            "",
            "Despite limitations, three-essay design + natural experiment + 27+ robustness tests"
        ]
    )

    # SLIDE 24: Conclusions
    add_content_slide(
        "Key Takeaways",
        [
            "1. MARKETS: Disclosure timing does NOT affect valuations; firm characteristics do",
            "2. UNCERTAINTY: Mandatory speed requirements INCREASE volatility through quality tradeoffs",
            "3. GOVERNANCE: Mandatory disclosure ACCELERATES governance response via stakeholder pressure",
            "4. POLICY: Current 7-day rule creates $0.76B costs without benefits; consider 45-60 day windows",
            "5. MECHANISMS: Three independent channels reveal regulation is more complex than \"faster = better\""
        ]
    )

    # SLIDE 25: Dashboard Preview
    add_content_slide(
        "Next: Interactive Dashboard",
        [
            "We will now walk through the interactive Streamlit dashboard showing:",
            "",
            "✓ Live regression results by essay and specification",
            "✓ Effect sizes with confidence intervals and significance stars",
            "✓ Size quartile heterogeneity analysis",
            "✓ Time series of breach frequency and average CAR by year",
            "✓ Economic significance scenarios by firm size",
            "✓ Validation test results (pre/post, industry controls, falsification)",
            "✓ Filter by FCC status, breach type, firm size"
        ]
    )

    # SLIDE 26: Final Slide
    add_title_slide(
        "Questions & Discussion",
        "Dissertation: Data Breach Disclosure Timing, Regulatory Burden, and Market Outcomes"
    )

    # Save presentation
    output_path = r'C:\Users\mcobp\BA798_TIM\Dissertation_Presentation.pptx'
    prs.save(output_path)
    print(f"Presentation created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == '__main__':
    create_presentation()
