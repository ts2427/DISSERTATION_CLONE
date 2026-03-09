"""
Create professional PowerPoint presentation with conceptual models
For committee and conference presentations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
COLORS = {
    'primary': RGBColor(31, 119, 180),      # Blue
    'accent': RGBColor(214, 39, 40),        # Red
    'success': RGBColor(44, 160, 44),       # Green
    'warning': RGBColor(255, 127, 14),      # Orange
    'dark': RGBColor(51, 51, 51),           # Dark gray
    'light': RGBColor(245, 245, 245)        # Light gray
}

def add_title_slide(prs, title, subtitle, author=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['primary']

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_p.alignment = PP_ALIGN.CENTER

    # Author
    if author:
        author_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        author_frame = author_box.text_frame
        author_p = author_frame.paragraphs[0]
        author_p.text = author
        author_p.font.size = Pt(16)
        author_p.font.color.rgb = RGBColor(200, 200, 200)
        author_p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, image_path=None, bullet_points=None, content_text=None):
    """Add a content slide with title, optional image, and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS['primary']
    title_shape.line.color.rgb = COLORS['primary']

    # Title text
    title_frame = title_shape.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.LEFT
    title_frame.margin_left = Inches(0.3)
    title_frame.margin_right = Inches(0.3)

    # Content
    if image_path and os.path.exists(image_path):
        # Add image
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.0), width=Inches(9))
    elif bullet_points:
        # Add bullet points
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.8))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        for i, bullet in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = bullet['text']
            p.level = bullet.get('level', 0)
            p.font.size = Pt(20 - (p.level * 3))
            p.font.color.rgb = COLORS['dark']
            p.space_before = Pt(6)
            p.space_after = Pt(6)

    elif content_text:
        # Add formatted text
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.8))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = content_text
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS['dark']
        p.line_spacing = 1.5

def add_two_column_slide(prs, title, left_image=None, right_image=None, left_text=None, right_text=None):
    """Add a two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS['primary']
    title_shape.line.color.rgb = COLORS['primary']

    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_frame.margin_left = Inches(0.3)

    # Left column
    if left_image and os.path.exists(left_image):
        slide.shapes.add_picture(left_image, Inches(0.3), Inches(1.0), width=Inches(4.7))
    elif left_text:
        left_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(4.7), Inches(6))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        p = left_frame.paragraphs[0]
        p.text = left_text
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark']

    # Right column
    if right_image and os.path.exists(right_image):
        slide.shapes.add_picture(right_image, Inches(5.0), Inches(1.0), width=Inches(4.7))
    elif right_text:
        right_box = slide.shapes.add_textbox(Inches(5.0), Inches(1.0), Inches(4.7), Inches(6))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        p = right_frame.paragraphs[0]
        p.text = right_text
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark']

# ============================================================================
# BUILD PRESENTATION
# ============================================================================

# Slide 1: Title
add_title_slide(prs,
    "Data Breach Disclosure Timing\nand Market Reactions",
    "How Regulatory Mandates Affect Markets, Information Asymmetry, and Governance",
    "Committee Presentation | 2025")

# Slide 2: Research Question & Motivation
add_content_slide(prs, "Central Research Question",
    bullet_points=[
        {'text': 'Is there any benefit to disclosing a data breach immediately?', 'level': 0},
        {'text': 'Or should disclosure timing be flexible?', 'level': 0},
        {'text': 'Regulatory Context: FCC Rule 37.3 (2007) mandates 7-day disclosure', 'level': 0},
        {'text': 'Data: 1,054 publicly-traded firm-breach observations, 2004-2025', 'level': 0},
        {'text': 'Gap: No empirical test of whether timing mandates improve outcomes', 'level': 0},
    ])

# Slide 3: Literature Foundation
add_content_slide(prs, "Theoretical Foundation",
    image_path='outputs/figures/CONCEPTUAL_01_LITERATURE_GENEALOGY.png')

# Slide 4: How Theories Build
add_content_slide(prs, "Information Asymmetry Theory: A Timeline",
    bullet_points=[
        {'text': 'Akerlof (1970): Bad information creates market failure', 'level': 0},
        {'text': 'Spence (1973): Quality revealed through costly signals', 'level': 0},
        {'text': 'Myers & Majluf (1984): Managers signal quality through timing', 'level': 0},
        {'text': 'Diamond & Verrecchia (1991): Forced disclosure can INCREASE asymmetry', 'level': 0},
        {'text': 'Tushman & Nadler (1978): Information processing capacity limits', 'level': 0},
        {'text': 'Your Research: Empirical test of regulatory timing effects', 'level': 0},
    ])

# Slide 5: The Complete Picture
add_content_slide(prs, "The Research Design: Natural Experiment",
    bullet_points=[
        {'text': 'Treatment: FCC-regulated telecommunications firms (N=200)', 'level': 0},
        {'text': 'Control: Non-FCC firms in same economy (N=854)', 'level': 0},
        {'text': 'Shock: FCC Rule 37.3 implemented January 1, 2007', 'level': 0},
        {'text': 'Strategy: Difference-in-differences with parallel trends validation', 'level': 0},
        {'text': 'Three Essays = Three Outcomes: Returns, Volatility, Governance', 'level': 0},
    ])

# Slide 6: Overarching Mechanism
add_content_slide(prs, "How It All Fits Together",
    image_path='outputs/figures/CONCEPTUAL_02_OVERARCHING_MECHANISM.png')

# Slide 7: Essay Models
add_content_slide(prs, "Three Essays, Three Mechanisms",
    image_path='outputs/figures/CONCEPTUAL_03_THREE_ESSAY_MODELS.png')

# Slide 8: Essay 1 Deep Dive
add_content_slide(prs, "ESSAY 1: Market Reactions (CAR)",
    bullet_points=[
        {'text': 'H1 (Timing): Does immediate disclosure reduce abnormal returns?', 'level': 0},
        {'text': '   Result: +0.57% (p=0.539) ❌ NOT significant', 'level': 1},
        {'text': 'H2 (FCC Effect): Do regulated firms experience worse market reactions?', 'level': 0},
        {'text': '   Result: -2.20%** (p=0.010) ✅ SIGNIFICANT & ROBUST', 'level': 1},
        {'text': 'H3 (Reputation): Do prior breaches predict worse reactions?', 'level': 0},
        {'text': '   Result: -0.08%*** per breach ✅ STRONGEST EFFECT', 'level': 1},
        {'text': 'Key Finding: Information environment (not timing) drives market pricing', 'level': 0},
    ])

# Slide 9: Essay 2 Deep Dive
add_content_slide(prs, "ESSAY 2: Information Asymmetry (Volatility)",
    bullet_points=[
        {'text': 'H5 (Pre-existing uncertainty): Does baseline volatility dominate?', 'level': 0},
        {'text': '   Result: Pre-volatility explains 68.6% of post-breach volatility ✅', 'level': 1},
        {'text': 'H2-Extended (FCC moderation): Do regulated firms see more uncertainty?', 'level': 0},
        {'text': '   Result: FCC increases volatility +1.83%** (p<0.05) ✅', 'level': 1},
        {'text': 'The Paradox: Forced disclosure INCREASES uncertainty (opposite of intent)', 'level': 0},
        {'text': 'Mechanism: Speed forces incompleteness; market prices it negatively', 'level': 0},
    ])

# Slide 10: Essay 3 Deep Dive
add_content_slide(prs, "ESSAY 3: Governance Response (Executive Turnover)",
    bullet_points=[
        {'text': 'H5 (Timing → Turnover): Does immediate disclosure trigger changes?', 'level': 0},
        {'text': '   Result: 46.4% executive turnover within 30 days (baseline)', 'level': 1},
        {'text': 'H6 (FCC moderation): Do regulatory firms see faster changes?', 'level': 0},
        {'text': '   Result: Modest FCC effect; governance response dominates', 'level': 1},
        {'text': 'Key Contrast: Governance response (46.4%) >> Regulatory enforcement (0.57%)', 'level': 0},
        {'text': 'Finding: Market discipline through governance > Regulatory punishment', 'level': 0},
    ])

# Slide 11: The Paradox
add_content_slide(prs, "⚠️ The Disclosure Paradox",
    bullet_points=[
        {'text': 'Regulator Logic (Sound): Faster disclosure → Less asymmetry → Better outcomes', 'level': 0},
        {'text': 'Implementation Problem: Timing mandate ≠ Information completeness', 'level': 0},
        {'text': 'Market Response: Incomplete disclosure signals bad news (adverse selection)', 'level': 0},
        {'text': 'Volatility Effect: Forced early disclosure INCREASES uncertainty', 'level': 0},
        {'text': 'Governance Effect: Stakeholders respond, but not to timing—to disclosure itself', 'level': 0},
        {'text': 'Bottom Line: Speed matters less than quality. Quality signals are credible.', 'level': 0},
    ])

# Slide 12: Policy Implications
add_content_slide(prs, "Policy Implications",
    image_path='outputs/figures/CONCEPTUAL_04_POLICY_OPTIONS.png')

# Slide 13: Policy Recommendations
add_content_slide(prs, "Three Evidence-Based Alternatives",
    bullet_points=[
        {'text': 'Option 1: Safe Harbor for Ongoing Investigations', 'level': 0},
        {'text': '   Allow delay when investigation is incomplete; mandatory final disclosure', 'level': 1},
        {'text': 'Option 2: Staged Disclosure (Preliminary + Final)', 'level': 0},
        {'text': '   Preliminary at 7 days + Complete within 30 days', 'level': 1},
        {'text': 'Option 3: Quality Standards, Not Speed Mandates', 'level': 0},
        {'text': '   Require completeness; let firms determine optimal timing', 'level': 1},
        {'text': 'Expected Outcome: Reduce market uncertainty while protecting investors', 'level': 0},
    ])

# Slide 14: Causal Identification & Robustness
add_content_slide(prs, "How We Know This is Causal (Not Correlation)",
    bullet_points=[
        {'text': 'Parallel Trends Test: FCC & non-FCC firms moved together pre-2007', 'level': 0},
        {'text': 'Post-2007 Emergence: FCC effect appears exactly when regulation takes effect', 'level': 0},
        {'text': 'Industry Controls: Effect strengthens with industry FE (rules out industry selection)', 'level': 0},
        {'text': 'Size Sensitivity: Effect stronger in small firms (not firm-size driven)', 'level': 0},
        {'text': 'Propensity Score Matching: Effect robust to selection on observables', 'level': 0},
        {'text': '25+ Robustness Specifications: All confirm main findings', 'level': 0},
    ])

# Slide 15: Data & Sample
add_content_slide(prs, "The Data: Comprehensive & Validated",
    bullet_points=[
        {'text': 'Full Sample: 1,054 publicly-traded firm-breach observations', 'level': 0},
        {'text': 'Study Period: 2004-2025 (16,104 firm-years)', 'level': 0},
        {'text': 'Essay 1 Sample: 898 breaches with CRSP stock data (87.8%)', 'level': 0},
        {'text': 'Essay 2 Sample: 891 breaches with volatility data (84.5%)', 'level': 0},
        {'text': 'Essay 3 Sample: 896 breaches with governance data (85.0%)', 'level': 0},
        {'text': 'Data Quality: 98% breach date accuracy (50 spot-check validation)', 'level': 0},
    ])

# Slide 16: Integrated Narrative
add_content_slide(prs, "Complete Research Narrative",
    image_path='outputs/figures/CONCEPTUAL_05_INTEGRATED_FLOW.png')

# Slide 17: Contributions
add_content_slide(prs, "Research Contributions",
    bullet_points=[
        {'text': 'Academic: First empirical test of regulatory timing using natural experiment', 'level': 0},
        {'text': 'Theory: Tests Myers & Majluf in mandatory disclosure context', 'level': 0},
        {'text': 'Method: Separates three mechanisms (valuation, uncertainty, governance)', 'level': 0},
        {'text': 'Data: Telecommunications sector first focus (previously understudied)', 'level': 0},
        {'text': 'Policy: Evidence that speed mandates can backfire if they force incompleteness', 'level': 0},
    ])

# Slide 18: Limitations & Future Research
add_content_slide(prs, "Limitations & Future Research",
    bullet_points=[
        {'text': 'Limitation 1: Sample limited to publicly-traded firms', 'level': 0},
        {'text': '   Future: Private firms may show different governance response', 'level': 1},
        {'text': 'Limitation 2: Causal identification relies on parallel trends assumption', 'level': 0},
        {'text': '   Future: Synthetic control methods could strengthen identification', 'level': 1},
        {'text': 'Limitation 3: 30/90-day executive turnover windows may miss long-run effects', 'level': 0},
        {'text': '   Future: Extended window analysis (1-3 years)', 'level': 1},
        {'text': 'Future Question: Do replacement executives improve security outcomes?', 'level': 0},
    ])

# Slide 19: Conclusion
add_content_slide(prs, "Key Takeaways",
    bullet_points=[
        {'text': "1. Timing regulations don't automatically improve outcomes", 'level': 0},
        {'text': '2. Market prices quality, not speed (Myers & Majluf confirmed)', 'level': 0},
        {'text': '3. Forced early disclosure can increase uncertainty (opposite of intent)', 'level': 0},
        {'text': '4. Governance response shows stakeholders care about disclosure, not timing', 'level': 0},
        {'text': '5. Policy should prioritize completeness over speed mandates', 'level': 0},
        {'text': 'Bottom Line: There is no substitute for good information.', 'level': 0},
    ])

# Slide 20: Thank You
add_title_slide(prs,
    "Questions?",
    "Data available at: BA798_TIM GitHub Repository\nPresentation & Analysis Scripts: Fully reproducible",
    "")

# Save presentation
output_path = 'Dissertation_Committee_Presentation.pptx'
prs.save(output_path)
print(f"[OK] Presentation saved: {output_path}")
