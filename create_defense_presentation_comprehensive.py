#!/usr/bin/env python3
"""
Create comprehensive dissertation defense presentation (22-25 slides)
Three essays: Market Returns, Information Asymmetry, Governance Response
Natural experiment: FCC Rule 37.3 (2007)
Sample: 1,054 breaches, 2006-2025
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
COLOR_DARK_BLUE = RGBColor(27, 94, 148)      # Deep blue
COLOR_TEAL = RGBColor(46, 204, 113)         # Teal - Essay 1
COLOR_ORANGE = RGBColor(230, 126, 34)       # Orange - Essay 2
COLOR_RED = RGBColor(192, 57, 43)           # Red - Essay 3
COLOR_GRAY = RGBColor(149, 165, 166)        # Gray - neutral
COLOR_LIGHT_BG = RGBColor(236, 240, 241)    # Light background

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(240, 240, 240)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, color=COLOR_DARK_BLUE):
    """Add content slide with colored header"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = color
    title_shape.line.color.rgb = color

    # Title text
    title_frame = title_shape.text_frame
    title_frame.vertical_anchor = 1  # Middle
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(10)
    p.space_after = Pt(10)

    return slide

def add_text_box(slide, left, top, width, height, text, size=14, bold=False, color=RGBColor(0, 0, 0)):
    """Add text box to slide"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box

def add_bullet_points(slide, left, top, width, height, bullets, size=13):
    """Add bullet points"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = frame.paragraphs[0]
        else:
            p = frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    return box

# ============================================================================
# SECTION 1: OPENING & CONTEXT (Slides 1-3)
# ============================================================================

# Slide 1: Title Slide
add_title_slide(prs,
    "Disclosure Timing and Market Reactions",
    "A Natural Experiment in Data Breach Regulation\nThree Essays on Valuation, Uncertainty, and Governance\n\n1,054 breaches | 2006-2025 | FCC Rule 37.3 Natural Experiment")

# Slide 2: The Puzzle
slide = add_content_slide(prs, "The Puzzle: Timing Paradox", COLOR_DARK_BLUE)
add_text_box(slide, 0.5, 1.2, 9, 0.6,
    "Do disclosure timing mandates improve market outcomes?",
    size=18, bold=True)

bullets = [
    "Essay 1: Does disclosure timing affect stock valuations? → NO (H1: +0.57%, p=0.539)",
    "Essay 2: Does disclosure timing affect market uncertainty? → YES (FCC: +1.83%**, increases volatility)",
    "Essay 3: Does disclosure timing trigger governance response? → YES (FCC: +5.3pp** executive turnover)",
    " ",
    "KEY INSIGHT: Timing doesn't change WHAT markets conclude, but it DOES change HOW markets learn and WHEN firms respond"
]
add_bullet_points(slide, 0.5, 2.0, 9, 5, bullets, size=13)

# Slide 3: Why This Matters
slide = add_content_slide(prs, "Policy Context: Regulatory Landscape", COLOR_DARK_BLUE)
bullets = [
    "FCC: 7-day disclosure mandate (Rule 37.3, effective 2007)",
    "SEC: 4-day cybersecurity disclosure rule (2023)",
    "HIPAA: 60-day notification requirement",
    "State laws: 30-90 days (California: 30 days)",
    " ",
    "THE GAP: No empirical evidence on whether timing mandates create benefits or costs",
    "YOUR CONTRIBUTION: First clean natural experiment isolating timing effects from firm quality signals",
    " ",
    "SAMPLE: 200 FCC firms (19%) vs 854 non-FCC firms (81%) in same sample → perfect comparison"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=12)

# ============================================================================
# SECTION 2: THEORY & FRAMEWORK (Slides 4-6)
# ============================================================================

# Slide 4: Theoretical Framework
slide = add_content_slide(prs, "Three Competing Theories", COLOR_DARK_BLUE)

# Theory 1
add_text_box(slide, 0.5, 1.2, 2.8, 0.4, "THEORY 1: Signaling Theory", size=12, bold=True, color=RGBColor(46, 204, 113))
add_text_box(slide, 0.5, 1.6, 2.8, 1.8,
    "Myers & Majluf (1984)\nFast disclosure = confidence\nSlow disclosure = bad news hoarding\nPrediction: Faster → Better returns\nYour evidence: REJECTED (p=0.539)",
    size=10)

# Theory 2
add_text_box(slide, 3.6, 1.2, 2.8, 0.4, "THEORY 2: Forced Disclosure Paradox", size=12, bold=True, color=RGBColor(230, 126, 34))
add_text_box(slide, 3.6, 1.6, 2.8, 1.8,
    "Diamond & Verrecchia (1991)\nMandatory disclosure can INCREASE uncertainty\nForced before investigation complete\nPrediction: Timing req → More uncertainty\nYour evidence: SUPPORTED (Essay 2)",
    size=10)

# Theory 3
add_text_box(slide, 6.7, 1.2, 2.8, 0.4, "THEORY 3: Stakeholder Theory", size=12, bold=True, color=RGBColor(192, 57, 43))
add_text_box(slide, 6.7, 1.6, 2.8, 1.8,
    "Freeman (1984)\nDisclosure activates stakeholders\nImmediate disclosure forces response\nPrediction: Faster → Gov't response\nYour evidence: SUPPORTED (Essay 3)",
    size=10)

# Integration
add_text_box(slide, 0.5, 3.6, 9, 3.5,
    "KEY FINDING: All three theories are partially correct. They address different outcomes through different mechanisms. The puzzle is resolved when you recognize that disclosure timing affects three independent outcomes (valuation level, learning speed, organizational response) not just one generic 'market outcome'.",
    size=12, bold=True, color=COLOR_DARK_BLUE)

# Slide 5: Natural Experiment Design
slide = add_content_slide(prs, "Natural Experiment: FCC Rule 37.3 (2007)", COLOR_DARK_BLUE)

bullets = [
    "THE SHOCK: FCC Rule 37.3 adopted September 2007 → effective September 30, 2007",
    " ",
    "TREATMENT GROUP (FCC-Regulated): SIC 4813 (wireline), 4899 (wireless), 4841 (cable) = 200 firms",
    "CONTROL GROUP (Non-FCC): All other industries = 854 firms",
    " ",
    "WHY IT'S CLEAN:",
    "  • Regulatory mandate, not market-driven",
    "  • Applies to whole industry (not individual firm choice)",
    "  • Exogenous shock in specific year",
    "  • Pre-treatment balance verified (all characteristics p>0.05)",
    " ",
    "VISUAL: FIGURE_PARALLEL_TRENDS shows pre/post divergence in CAR trends"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=12)

# Slide 6: Identification Strategy
slide = add_content_slide(prs, "Four Robustness Tests: Causal Identification", COLOR_DARK_BLUE)

bullets = [
    "TEST 1 - Temporal Validation (Pre-2007):",
    "  Pre-2007 FCC coef: -13.96% (p=0.88, NS) | Post-2007: -2.26% (p=0.0125, sig)",
    "  ✓ Effect only appears after regulation, not before",
    " ",
    "TEST 2 - Industry Fixed Effects:",
    "  Baseline FCC: -2.20% | With industry FE: -5.37% (STRONGER!)",
    "  ✓ Not driven by industry composition, effect strengthens with controls",
    " ",
    "TEST 3 - Size Sensitivity (Quartile Analysis):",
    "  Shows effects vary by size but all in expected direction",
    "  ✓ Acknowledges confound but pattern is consistent",
    " ",
    "TEST 4 - Multi-Outcome Consistency:",
    "  Same natural experiment explains effects on returns, volatility, AND turnover",
    "  ✓ Not spurious; regulatory change affects multiple mechanisms"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=11)

# ============================================================================
# SECTION 3: ESSAY 1 - MARKET REACTIONS (Slides 7-10)
# ============================================================================

# Slide 7: Essay 1 Main Findings
slide = add_content_slide(prs, "Essay 1: Market Returns - Main Findings", COLOR_TEAL)

add_text_box(slide, 0.5, 1.2, 9, 0.5, "H1 - Timing Effect (Null Result):", size=13, bold=True, color=COLOR_TEAL)
add_text_box(slide, 0.7, 1.7, 8.6, 1.2,
    "Coefficient: +0.57% (p=0.539, NOT significant) | 90% CI: [-0.95%, +2.09%] | TOST Equivalence: PASS | Sample: 898 breaches",
    size=12)

add_text_box(slide, 0.5, 3.0, 9, 0.5, "What DOES Matter:", size=13, bold=True, color=COLOR_TEAL)

# Create simple table
table_data = [
    ["Factor", "Effect", "Significance", "Interpretation"],
    ["FCC Regulation", "-2.20%", "p=0.010**", "Regulatory burden/complexity penalty"],
    ["Health Breach", "-2.51%", "p=0.004***", "HIPAA compliance risk"],
    ["Prior Breaches", "-0.22% per breach", "p<0.001***", "STRONGEST: Market prices reputation"]
]

# Add table as text
add_text_box(slide, 0.7, 3.5, 8.6, 0.35, "FCC Regulation      -2.20%**     Regulatory burden", size=11, bold=True)
add_text_box(slide, 0.7, 3.95, 8.6, 0.35, "Health Breach        -2.51%***    HIPAA compliance risk", size=11, bold=True)
add_text_box(slide, 0.7, 4.4, 8.6, 0.35, "Prior Breaches       -0.22%***    STRONGEST: Market prices reputation", size=11, bold=True)

add_text_box(slide, 0.5, 5.0, 9, 2,
    "CENTRAL FINDING: Markets punish WHO YOU ARE and WHAT WAS BREACHED — not WHEN YOU TALK.\n\nThe null finding is meaningful: It directly contradicts regulatory assumptions that mandatory speed requirements create market benefits.",
    size=12, bold=True, color=RGBColor(192, 57, 43))

# Slide 8: Essay 1 - H1 Timing Explained
slide = add_content_slide(prs, "Essay 1: Why H1 Null is Robust", COLOR_TEAL)

bullets = [
    "NOT due to low power:",
    "  • Post-hoc power analysis: MDE = ±2.39pp (can detect economically meaningful effects)",
    "  • Observed effect (+0.57pp) much smaller than detectable threshold",
    "  • This SUPPORTS true null, not just under-powered study",
    " ",
    "Equivalence testing confirms:",
    "  • 90% Confidence Interval: [-0.95%, +2.09%]",
    "  • Falls entirely within ±2.10% equivalence bounds → effect is negligible",
    " ",
    "Robust across ALL specifications:",
    "  • 4 event windows × 7 timing thresholds × 8 subsamples × 6 SE methods = 27+ tests",
    "  • Timing effect non-significant across ALL"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=12)

# Slide 9: Essay 1 - Heterogeneity by Size
slide = add_content_slide(prs, "Essay 1: Firm Size Heterogeneity", COLOR_TEAL)

bullets = [
    "The Size Confound: FCC firms are 2x larger ($62.6B vs $31.0B, p<0.0001)",
    " ",
    "FCC Effect by Firm Size Quartile:",
    "  Q1 (Smallest): -6.22% (p=0.053, marginal)",
    "  Q2: -4.06% (p=0.007**)",
    "  Q3: +0.66% (p=0.703, NS)",
    "  Q4 (Largest): +0.43% (p=0.692, NS)",
    " ",
    "Interpretation: FCC penalty concentrated in smaller firms—they feel burden most acutely",
    " ",
    "H1 Timing Effect by Size: ALL non-significant",
    "  Q1: +1.158% | Q2: +0.847% | Q3: +0.123% | Q4: +0.050% (all p>0.05)",
    " ",
    "CONCLUSION: Timing irrelevant regardless of firm size (universal null finding)"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=11)

# Slide 10: Essay 1 - Robustness Summary
slide = add_content_slide(prs, "Essay 1: Robustness Across 27+ Specifications", COLOR_TEAL)

bullets = [
    "H1 Timing Tested Across Multiple Dimensions:",
    " ",
    "EVENT WINDOWS: 5d, 10d, 30d, 60d → All NS",
    "TIMING THRESHOLDS: 1d, 3d, 7d, 14d, 30d, 60d, 90d → All NS",
    "SUBSAMPLES: Full, no-prior, repeat, health, financial, FCC, non-FCC, pre-2007 → All NS",
    "SE METHODS: OLS, HC3, Clustered, Panel-corrected, Robust, Bootstrap → All NS",
    " ",
    "Controls tested: Firm financials, industry FE, year FE, firm FE → H1 null persists",
    "ML validation: Random Forest & Gradient Boosting → Timing not predictive",
    " ",
    "Output files:",
    "  • H1_TOST_Equivalence_Test.txt (power analysis)",
    "  • H1_Comprehensive_Power_Analysis.csv (MDE calculations)"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=11)

# ============================================================================
# SECTION 4: ESSAY 2 - INFORMATION ASYMMETRY (Slides 11-14)
# ============================================================================

# Slide 11: Essay 2 Paradoxical Finding
slide = add_content_slide(prs, "Essay 2: Information Asymmetry - Paradoxical Finding", COLOR_ORANGE)

add_text_box(slide, 0.5, 1.2, 9, 0.6, "The Paradox: Same regulation, opposite effects on different outcomes", size=13, bold=True)

bullets = [
    "RETURNS (Essay 1): FCC effect = -2.20% (bad, but doesn't change with timing)",
    "VOLATILITY (Essay 2): FCC effect = +1.83%** (undesirable uncertainty INCREASES)",
    " ",
    "Why is this interesting?",
    "  • NOT contradictory—reveals different mechanisms",
    "  • Regulatory constraint (timing requirement) creates information quality problem",
    "  • Forces disclosure BEFORE investigation complete",
    "  • Markets learn faster but with INCOMPLETE information = higher uncertainty",
    " ",
    "Sample: 891 breaches with volatility data",
    " ",
    "KEY INSIGHT: FCC creates market penalty on returns AND increases market uncertainty—simultaneous effects through different mechanisms"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=12)

# Slide 12: Essay 2 Mechanism
slide = add_content_slide(prs, "Essay 2: Timing-Quality Tradeoff Mechanism", COLOR_ORANGE)

bullets = [
    "THE MECHANISM (Timing-Quality Paradox):",
    " ",
    "1. FCC Rule 37.3 requires disclosure within 30 days",
    "2. Regulation creates time pressure on investigation completion",
    "3. Firms forced to disclose with INCOMPLETE breach analysis",
    "4. Markets see disclosure but recognize information is INCOMPLETE",
    "5. Result: HIGHER UNCERTAINTY (volatility increases)",
    " ",
    "Evidence for this mechanism:",
    "  • Effect strongest in SMALL firms (+7.31%*** in Q1)",
    "  • Weak in LARGE firms (-3.39%** in Q4, reverses sign!)",
    "  • Interpretation: Small firms have limited breach investigation capacity",
    "  • Large firms can conduct parallel processes with no quality loss",
    " ",
    "Timing coefficient: Non-significant for volatility",
    "  → Not about delay duration, but about regulatory CONSTRAINT on investigation"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=11)

# Slide 13: Essay 2 Firm Size Heterogeneity
slide = add_content_slide(prs, "Essay 2: THE CRITICAL FINDING - Firm Size Heterogeneity", COLOR_ORANGE)

add_text_box(slide, 0.5, 1.2, 9, 0.4, "FCC Effect on Volatility by Firm Size Quartile:", size=13, bold=True, color=COLOR_ORANGE)

# Quartile data
add_text_box(slide, 0.7, 1.7, 8.6, 0.3, "Q1 (Smallest):    +7.31%***    p<0.001    232 firms", size=11, bold=True, color=RGBColor(192, 57, 43))
add_text_box(slide, 0.7, 2.05, 8.6, 0.3, "Q2 (Small):       +3.64%**     p=0.015    233 firms", size=11, bold=True, color=RGBColor(230, 126, 34))
add_text_box(slide, 0.7, 2.4, 8.6, 0.3, "Q3 (Medium):      -0.54%       p=0.770    213 firms", size=11, bold=True, color=RGBColor(149, 165, 166))
add_text_box(slide, 0.7, 2.75, 8.6, 0.3, "Q4 (Largest):     -3.39%**     p=0.024    213 firms", size=11, bold=True, color=RGBColor(46, 204, 113))

add_text_box(slide, 0.5, 3.2, 9, 0.3, "11 PERCENTAGE POINT DIFFERENTIAL (Q1 vs Q4)", size=12, bold=True, color=COLOR_ORANGE)

bullets = [
    "The Pattern:",
    "  • Linear decline as firm size increases",
    "  • Sign reversal at large firms: FCC actually REDUCES volatility",
    "  • Small firms: Uncertainty increases (capacity constraint to handle timing pressure)",
    "  • Large firms: Volatility decreases (can handle mandatory timing without information loss)",
    " ",
    "Why this matters for policy:",
    "  • Regulation's costs are REGRESSIVE—smallest firms hurt most",
    "  • Economic impact: Small firms face ~$37M annual cost of capital increase per breach",
    "  • Largest firms have infrastructure neutralizing disclosure timing pressure"
]
add_bullet_points(slide, 0.5, 3.6, 9, 3.8, bullets, size=11)

# Slide 14: Essay 2 Alternative Mechanisms Tested
slide = add_content_slide(prs, "Essay 2: Alternative Mechanisms Tested & Rejected", COLOR_ORANGE)

add_text_box(slide, 0.5, 1.2, 9, 0.4, "Which mechanisms might explain the FCC volatility effect?", size=13, bold=True)

mechanisms = [
    ["Mechanism", "Result", "P-value", "Conclusion"],
    ["CVSS Technical Complexity", "-0.0784pp", "p=0.97", "NOT sig—technical attributes don't drive"],
    ["Governance Quality", "+1.6989", "p=0.84", "NOT sig—firm governance doesn't moderate"],
    ["Information Environment", "-2.6142", "p=0.28", "NOT sig—media coverage doesn't explain"],
    ["Complexity Index (Script 105)", "-0.0784pp", "p=0.97", "NOT sig—technical complexity irrelevant"],
    ["Media Coverage (Spec A)", "+0.5585", "p=0.034", "Sig but small—not primary driver"]
]

for i, row in enumerate(mechanisms):
    y_pos = 1.6 + (i * 0.34)
    if i == 0:  # Header
        add_text_box(slide, 0.7, y_pos, 1.8, 0.3, row[0], size=10, bold=True)
        add_text_box(slide, 2.6, y_pos, 1.5, 0.3, row[1], size=10, bold=True)
        add_text_box(slide, 4.2, y_pos, 1.2, 0.3, row[2], size=10, bold=True)
        add_text_box(slide, 5.5, y_pos, 3.8, 0.3, row[3], size=10, bold=True)
    else:
        color = RGBColor(192, 57, 43) if "NOT sig" in row[2] else RGBColor(46, 204, 113)
        add_text_box(slide, 0.7, y_pos, 1.8, 0.3, row[0], size=9, color=color)
        add_text_box(slide, 2.6, y_pos, 1.5, 0.3, row[1], size=9, color=color)
        add_text_box(slide, 4.2, y_pos, 1.2, 0.3, row[2], size=9, color=color)
        add_text_box(slide, 5.5, y_pos, 3.8, 0.3, row[3], size=9, color=color)

add_text_box(slide, 0.5, 4.0, 9, 3,
    "KEY FINDING: Of all possible mechanisms, ONLY FIRM SIZE matters.\n\nThis is the critical heterogeneity result—it's not driven by breach characteristics (CVSS), but by firm organizational CAPACITY to handle regulatory timing constraints.\n\nThis isolates the mechanism: It's a capacity constraint problem, not a breach severity problem.",
    size=12, bold=True, color=RGBColor(192, 57, 43))

# ============================================================================
# SECTION 5: CRITICAL INSIGHT (Slide 15)
# ============================================================================

# Slide 15: CRITICAL INTEGRATION - Essay 2 vs Essay 3 Tension
slide = add_content_slide(prs, "CRITICAL INSIGHT: Timing Affects Learning Speed, Not Valuation Level", COLOR_DARK_BLUE)

add_text_box(slide, 0.5, 1.2, 9, 0.5, "The Apparent Contradiction Resolved:", size=13, bold=True)

# Three-column integration
add_text_box(slide, 0.5, 1.8, 2.8, 0.3, "ESSAY 1: Returns", size=11, bold=True, color=COLOR_TEAL)
add_text_box(slide, 0.5, 2.1, 2.8, 2.5,
    "Effect: NO EFFECT\n\nMechanism: Timing doesn't change what markets conclude\n\nSpeed Effect: Timing IRRELEVANT to valuations\n\nFinding: Market reaches same conclusion whether fast or slow",
    size=10)

add_text_box(slide, 3.6, 1.8, 2.8, 0.3, "ESSAY 2: Volatility", size=11, bold=True, color=COLOR_ORANGE)
add_text_box(slide, 3.6, 2.1, 2.8, 2.5,
    "Effect: INCREASES +1.83%**\n\nMechanism: Forced incomplete disclosure\n\nSpeed Effect: Speed creates QUALITY LOSS\n\nFinding: Uncertainty SPIKE when forced to disclose early",
    size=10)

add_text_box(slide, 6.7, 1.8, 2.8, 0.3, "ESSAY 3: Turnover", size=11, bold=True, color=COLOR_RED)
add_text_box(slide, 6.7, 2.1, 2.8, 2.5,
    "Effect: ACCELERATES +5.3pp**\n\nMechanism: Stakeholder activation\n\nSpeed Effect: Speed TRIGGERS RESPONSE\n\nFinding: Organizational response FASTER with public disclosure",
    size=10)

# Key insight box
add_text_box(slide, 0.5, 4.8, 9, 2.5,
    "THE KEY DISTINCTION: Timing affects LEARNING SPEED, not VALUATION LEVEL\n\n• Markets learn equally fast or slow (doesn't change CAR) ✓ Essay 1 null\n• BUT learning process matters—forced speed creates incomplete information ✓ Essay 2 finding\n• AND forced disclosure activates stakeholders who demand organizational response ✓ Essay 3 finding\n\nAll three are operating SIMULTANEOUSLY through DIFFERENT channels.",
    size=12, bold=True, color=RGBColor(192, 57, 43))

# ============================================================================
# SECTION 6: ESSAY 3 - GOVERNANCE RESPONSE (Slides 16-18)
# ============================================================================

# Slide 16: Essay 3 Main Findings
slide = add_content_slide(prs, "Essay 3: Governance Response - Executive Turnover", COLOR_RED)

bullets = [
    "Research Question: Does disclosure timing trigger executive governance changes?",
    " ",
    "THE FINDINGS:",
    "  • Baseline executive turnover (all breaches, 30-day window): 46.4%",
    "  • With immediate disclosure: 50.6%",
    "  • With delayed disclosure: 45.3%",
    "  • DIFFERENCE: +5.3 percentage points (p=0.008**)",
    " ",
    "Effect Dynamics:",
    "  • PEAK turnover at 30 days post-disclosure",
    "  • Effect persists but decays: Still significant at 90 days (66.9% cumulative)",
    "  • By 180 days: Stabilizes (67.5% cumulative)",
    "  • INTERPRETATION: Acute crisis response, not embedded governance reform",
    " ",
    "Sample: 896 breaches | Scope: Officers & directors | Mean: 3.2 executives per breach"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=11)

# Slide 17: Essay 3 Mechanism
slide = add_content_slide(prs, "Essay 3: Stakeholder Activation Mechanism", COLOR_RED)

bullets = [
    "WHY DOES IMMEDIATE DISCLOSURE TRIGGER TURNOVER?",
    " ",
    "Step 1: Disclosure Act",
    "  → Public announcement of breach activates multiple stakeholders",
    " ",
    "Step 2: Stakeholder Pressure",
    "  → Investors, employees, customers, regulators now aware and concerned",
    " ",
    "Step 3: Board Response",
    "  → Board sees breach as organizational crisis requiring visible governance action",
    " ",
    "Step 4: Executive Turnover",
    "  → Changes signal accountability, governance competence, fresh oversight",
    " ",
    "CRITICAL VALIDATION: Volatility does NOT mediate timing → turnover",
    "  • Mediation analysis (Script 91) shows indirect effect through volatility = -0.0114 (NS)",
    "  • 95% CI includes zero: [-0.0433, 0.0205]",
    "  • CONCLUSION: Governance response is DIRECT stakeholder pressure, not information-driven"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=11)

# Slide 18: Essay 3 Size Sensitivity
slide = add_content_slide(prs, "Essay 3: Size Heterogeneity & Robustness", COLOR_RED)

add_text_box(slide, 0.5, 1.2, 9, 0.4, "Size Heterogeneity (Timing Effect on Turnover by Quartile):", size=12, bold=True)

add_text_box(slide, 0.7, 1.65, 8.6, 0.25, "Q1 (Smallest):    -0.679     p=0.081 ~    Marginal (negative)", size=10)
add_text_box(slide, 0.7, 1.95, 8.6, 0.25, "Q2 (Med-small):  -1.132**   p=0.026 **   Stronger response", size=10)
add_text_box(slide, 0.7, 2.25, 8.6, 0.25, "Q3 (Med-large):  -1.651***  p=0.006 ***  STRONGEST response", size=10)
add_text_box(slide, 0.7, 2.55, 8.6, 0.25, "Q4 (Largest):     +0.371     p=0.265      NO response", size=10)

bullets = [
    " ",
    "Pattern Interpretation:",
    "  • Medium firms (Q2/Q3) show STRONGEST governance response",
    "  • Smallest firms show marginal effect (hard to change structure with limited boards)",
    "  • Largest firms show NO effect (sophisticated IR practices shield from stakeholder pressure)",
    " ",
    "NOTE: Size pattern is DIFFERENT from volatility (Essay 2)",
    "  • Essay 2: Small firms hurt (+7.31%), large firms helped (-3.39%)",
    "  • Essay 3: Medium firms respond most (logit coefficients)",
    "  • VALIDATES: Separate mechanisms, not spillover effects"
]
add_bullet_points(slide, 0.5, 3.0, 9, 4.4, bullets, size=11)

# ============================================================================
# SECTION 7: SYNTHESIS & POLICY (Slides 19-21)
# ============================================================================

# Slide 19: Integration Table
slide = add_content_slide(prs, "Cross-Essay Integration: Three Independent Mechanisms", COLOR_DARK_BLUE)

add_text_box(slide, 0.5, 1.2, 9, 0.35, "How Disclosure Requirements Work Through Multiple Pathways:", size=12, bold=True)

# Simplified table
add_text_box(slide, 0.5, 1.65, 1.5, 0.3, "Essay 1", size=10, bold=True, color=COLOR_TEAL)
add_text_box(slide, 0.5, 2.0, 1.5, 2.3,
    "Outcome: Stock Returns\n\nEffect: NO effect\n\nMechanism: Timing irrelevant\n\nPath: Valuation Level\n\n(UNCHANGED)",
    size=9)

add_text_box(slide, 2.15, 1.65, 1.5, 0.3, "Essay 2", size=10, bold=True, color=COLOR_ORANGE)
add_text_box(slide, 2.15, 2.0, 1.5, 2.3,
    "Outcome: Volatility\n\nEffect: +1.83%**\n\nMechanism: Incomplete disclosure\n\nPath: Information Quality\n\n(DECREASED)",
    size=9)

add_text_box(slide, 3.8, 1.65, 1.5, 0.3, "Essay 3", size=10, bold=True, color=COLOR_RED)
add_text_box(slide, 3.8, 2.0, 1.5, 2.3,
    "Outcome: Turnover\n\nEffect: +5.3pp**\n\nMechanism: Stakeholder pressure\n\nPath: Org. Response\n\n(ACCELERATED)",
    size=9)

add_text_box(slide, 0.5, 4.4, 9, 3,
    "CENTRAL FINDING: Same regulation (FCC Rule 37.3) produces three distinct outcomes through three separate mechanisms. All are operating SIMULTANEOUSLY. The regulation achieves some goals (governance response) while creating unintended consequences (information quality loss) and leaving others unaffected (market valuations).\n\nThis resolves the apparent contradictions and explains policy complexity.",
    size=11, bold=True, color=RGBColor(192, 57, 43))

# Slide 20: Policy Implications
slide = add_content_slide(prs, "Policy Implications: Tightened & Evidence-Based", COLOR_DARK_BLUE)

add_text_box(slide, 0.5, 1.2, 9, 0.35, "Four Evidence-Based Claims (with proper scope):", size=12, bold=True)

bullets = [
    "CLAIM 1: Stock Market Discipline Does NOT Operate Through Timing",
    "  Evidence: H1 null (+0.57%, p=0.539), robust across 27+ specs, TOST confirmed",
    "  Implication: Market does not penalize slow disclosure or reward fast disclosure",
    "  SCOPE: Stock market shareholder reactions ONLY (not consumer protection, compliance, trust)",
    "  Why: Explains puzzle—if markets rewarded speed, firms would volunteer. They don't.",
    " ",
    "CLAIM 2: Timing Requirements Create Information Quality Tradeoffs",
    "  Evidence: Essay 2 (+1.83%** volatility), mechanism = forced incomplete disclosure",
    "  Implication: Regulatory timing constraints increase market uncertainty (unintended consequence)",
    "  SCOPE: FCC-regulated telecommunications firms; effect strongest in small firms",
    "  Why: Regulations achieve disclosure SPEED but at cost of completeness",
    " ",
    "CLAIM 3: Disclosure Successfully Activates Governance Response",
    "  Evidence: Essay 3 (+5.3pp** turnover), NOT mediated by information quality",
    "  Implication: Disclosure effectively drives organizational accountability through stakeholder pressure",
    "  SCOPE: Executive/board-level changes; effect strongest in medium-sized firms",
    "  Why: Regulation's primary organizational control mechanism IS working"
]
add_bullet_points(slide, 0.5, 1.6, 9, 5.8, bullets, size=10)

# Slide 21: Limitations & Future Research
slide = add_content_slide(prs, "Limitations & Future Research Agenda", COLOR_DARK_BLUE)

bullets = [
    "KEY LIMITATIONS:",
    "  ✓ Scope: Stock market reactions only (don't observe consumer harm, compliance costs, trust)",
    "  ✓ Sample: Publicly-traded firms only (private firm dynamics may differ)",
    "  ✓ Time period: 2006-2025 (telecom industry evolving)",
    "  ✓ Windows: 30/90/180 days (longer-term effects beyond scope)",
    "  ✓ Causality: Natural experiment reduces but doesn't eliminate endogeneity",
    " ",
    "FUTURE RESEARCH:",
    "  1. Cross-sector testing: Do mechanisms replicate for SEC cybersecurity rule? HIPAA?",
    "  2. Consumer outcomes: Link stock reactions to actual consumer harm measures",
    "  3. Regulatory compliance: Do breach investigation practices change post-disclosure?",
    "  4. International: How do other countries balance speed vs. completeness? (GDPR)",
    "  5. Mechanism validation: Can we directly measure 'investigation completeness'?",
    "  6. Long-term governance: Do executive changes persist beyond 180 days?",
    " ",
    "CLOSING: Disclosure policy is complex. Single-dimension theories (timing) fail.",
    "Move beyond 'faster is better' and explicitly design regulations balancing competing goals."
]
add_bullet_points(slide, 0.5, 1.2, 9, 6.2, bullets, size=10)

# ============================================================================
# SECTION 8: APPENDIX (Optional - Slides 22-25)
# ============================================================================

# Slide 22: Robustness Summary
slide = add_content_slide(prs, "Appendix: Robustness Testing Matrix", COLOR_GRAY)

bullets = [
    "H1 Timing Effect Tested Across 27+ Specifications:",
    " ",
    "DIMENSIONS TESTED:",
    "  • Event Windows: 5d, 10d, 30d, 60d (all NS)",
    "  • Timing Thresholds: 1d, 3d, 7d, 14d, 30d, 60d, 90d (all NS)",
    "  • Subsamples: Full, no-prior, repeat, health, financial, FCC, non-FCC, pre-2007 (all NS)",
    "  • Standard Errors: OLS, HC3, Clustered, Panel-corrected, Robust, Bootstrap (all NS)",
    "  • Controls: Firm financials, Industry FE, Year FE, Firm FE (H1 null persists)",
    " ",
    "VALIDATION:",
    "  • ML Models: Random Forest & Gradient Boosting (timing not predictive)",
    "  • Parallel Trends: Pre/post divergence visible only post-2007 (supports identification)",
    "  • Balance Test: All characteristics balanced pre-2007, p>0.05 (no pre-treatment confound)"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=11)

# Slide 23: Data Quality
slide = add_content_slide(prs, "Appendix: Data Quality & Coverage", COLOR_GRAY)

bullets = [
    "MATCHING RATES:",
    "  • CRSP match: 92.1% of raw breach records (926 of 1,006) → 87.9% of final dataset (926 of 1,054)",
    "  • Compustat link: 100% (CIK matching)",
    "  • CVSS scores: 80.4% (847 of 1,054 with CVE data)",
    "  • SEC 8-K data: 90.5% (956 of 1,054 with filings)",
    " ",
    "SAMPLE FLOW:",
    "  1,054 total breaches → 926 with CRSP data → 898 with complete Essay 1 data",
    "  1,054 total breaches → 916 with volatility data → 891 complete Essay 2",
    "  1,054 total breaches → 896 with executive change data → Essays 3",
    " ",
    "MISSING DATA POLICY:",
    "  • Complete-case analysis for each essay (no imputation)",
    "  • Different N by essay reflects different data requirements, not systematic selection"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=11)

# Slide 24: Economic Significance
slide = add_content_slide(prs, "Appendix: Economic Significance", COLOR_GRAY)

bullets = [
    "AGGREGATE ECONOMIC IMPACT:",
    "  • Empirical sample aggregate: $0.76B (total shareholder losses across 1,054 breaches)",
    "  • Theoretical annual projection: $9.9B (if breach incident rate continues)",
    " ",
    "FIRM-LEVEL IMPACTS:",
    "  • Small firms (Essay 2): ~$37M additional annual cost of capital per breach",
    "    (from +7.31% volatility increase × typical firm market cap)",
    " ",
    "REGULATORY COSTS:",
    "  • FCC compliance documentation & investigation pressure",
    "  • Information quality loss (forced incomplete disclosure)",
    "  • Executive turnover costs (lost institutional knowledge)",
    " ",
    "REGULATORY BENEFITS:",
    "  • Governance response: 50.6% vs 45.3% turnover (boards do respond to pressure)",
    "  • Stakeholder accountability activation",
    "  • Public information disclosure (market learning, though with increased uncertainty)"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=11)

# Slide 25: Key Papers & Methodological References
slide = add_content_slide(prs, "Appendix: Theoretical Foundation & Methods", COLOR_GRAY)

bullets = [
    "THEORETICAL PILLARS:",
    "  • Myers & Majluf (1984): Signaling Theory (tested & rejected for timing)",
    "  • Diamond & Verrecchia (1991): Forced Disclosure Paradox (supported for volatility)",
    "  • Freeman (1984): Stakeholder Theory (supported for governance response)",
    " ",
    "METHODOLOGICAL FOUNDATIONS:",
    "  • Brown & Warner (1985): Event study methodology (CAR, event windows)",
    "  • Fama & French (1993): Three-factor model (abnormal returns)",
    "  • Natural experiments (Angrist & Pischke 2009): Causal identification",
    " ",
    "SPECIFIC METHODS USED:",
    "  • OLS with HC3 robust standard errors (Essays 1-2)",
    "  • Logit with cluster SEs (Essay 3)",
    "  • Mediation analysis with delta-method SEs (Script 91)",
    "  • TOST Equivalence testing (H1 power analysis)",
    "  • Parallel trends visualization (causal identification)",
    "  • Balance testing (pre-treatment confound check)"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=10)

# ============================================================================
# Save Presentation
# ============================================================================

output_path = 'Dissertation_Defense_Presentation.pptx'
prs.save(output_path)

print("=" * 80)
print("SUCCESS: Dissertation Defense Presentation Created")
print("=" * 80)
print(f"\nFile: {output_path}")
print(f"Total Slides: {len(prs.slides)}")
print("\nSLIDE BREAKDOWN:")
print("  Section 1 (Opening): Slides 1-3")
print("  Section 2 (Theory): Slides 4-6")
print("  Section 3 (Essay 1): Slides 7-10")
print("  Section 4 (Essay 2): Slides 11-14")
print("  Section 5 (CRITICAL INSIGHT): Slide 15")
print("  Section 6 (Essay 3): Slides 16-18")
print("  Section 7 (Integration & Policy): Slides 19-21")
print("  Section 8 (Appendix): Slides 22-25")
print("\nKEY FEATURES:")
print("  [OK] Expert feedback incorporated (Essay 2 vs 3 tension on Slide 15)")
print("  [OK] Tightened policy implications with proper scope bounds")
print("  [OK] All heterogeneity findings documented")
print("  [OK] Causal identification thoroughly explained")
print("  [OK] Robustness tests visible")
print("  [OK] Data-backed claims with exact numbers")
print("  [OK] 22-25 slides optimized for 35-40 minute defense + questions")
print("\nREADY FOR APRIL 23, 2026 DEFENSE")
print("=" * 80)
