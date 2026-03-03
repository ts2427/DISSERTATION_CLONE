"""
Update Dissertation_Presentation_Final.pptx with correct H1 values
"""

from pptx import Presentation
import re

pptx_path = r'C:\Users\mcobp\DISSERTATION_CLONE\Dissertation_Presentation_Final.pptx'

# Load presentation
prs = Presentation(pptx_path)

# Patterns to replace
replacements = [
    # H1 coefficient
    (r'\+0\.57%', '+0.649%'),
    (r'0\.57%', '0.649%'),
    (r'\+0\.57pp', '+0.649pp'),

    # H1 p-value
    ('p=0.373', 'p=0.443'),
    ('p = 0.373', 'p = 0.443'),
    ('p=0.539', 'p=0.443'),
    ('p = 0.539', 'p = 0.443'),
]

changes_made = 0
slides_modified = set()

# Process all slides
for slide_num, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            original_text = shape.text
            new_text = original_text

            # Apply replacements
            for pattern, replacement in replacements:
                if re.search(pattern, new_text):
                    new_text = re.sub(pattern, replacement, new_text)

            # If text changed, update it
            if new_text != original_text:
                if hasattr(shape, 'text_frame'):
                    # Clear existing text and add new
                    shape.text_frame.clear()

                    # Split by lines to preserve structure
                    for line in new_text.split('\n'):
                        p = shape.text_frame.add_paragraph()
                        p.text = line

                    changes_made += 1
                    slides_modified.add(slide_num)
                    print(f"Slide {slide_num}: Updated")
                    print(f"  From: {original_text[:60]}...")
                    print(f"  To:   {new_text[:60]}...")

# Save updated presentation
prs.save(pptx_path)

print()
print("=" * 80)
print(f"Presentation Updated: {changes_made} text elements modified")
print(f"Slides affected: {sorted(slides_modified)}")
print("=" * 80)
print()
print("Updated Values:")
print("  - H1 coefficient: +0.649% (was 0.57%)")
print("  - H1 p-value: 0.443 (was 0.373/0.539)")
print()
print(f"File saved: {pptx_path}")
