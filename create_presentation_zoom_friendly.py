"""
Zoom-Friendly Dissertation Presentation (PowerPoint)
Larger text, less density, more slides, better readability
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_presentation():
    """Generate the Zoom-friendly dissertation presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Color scheme
    DARK_BLUE = RGBColor(31, 78, 121)
    LIGHT_BLUE = RGBColor(79, 129, 189)
    ACCENT_RED = RGBColor(192, 0, 0)
    ACCENT_GREEN = RGBColor(0, 176, 80)
    WHITE = RGBColor(255, 255, 255)
    DARK_GRAY = RGBColor(89, 89, 89)

    def add_title_slide(title, subtitle=""):
        """Add a title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BLUE

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            p = subtitle_frame.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = LIGHT_BLUE

        return slide

    def add_content_slide(title, bullet_points=None):
        """Add a content slide with title and bullets (ZOOM-FRIENDLY: fewer, bigger bullets)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title box
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = DARK_BLUE
        title_shape.line.color.rgb = DARK_BLUE

        title_frame = title_shape.text_frame
        title_frame.word_wrap = False
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Content box
        if bullet_points:
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True

            for i, point in enumerate(bullet_points):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                p.text = point
                p.font.size = Pt(24)  # Larger for Zoom
                p.font.color.rgb = DARK_GRAY
                p.space_before = Pt(16)
                p.space_after = Pt(16)
                p.level = 0

        return slide

    def add_key_metric_slide(title, metric_title, metric_value, metric_unit="", description=""):
        """Add a slide highlighting a key metric (ZOOM-FRIENDLY)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title box
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = DARK_BLUE
        title_shape.line.color.rgb = DARK_BLUE

        title_frame = title_shape.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Large metric box
        metric_shape = slide.shapes.add_shape(1, Inches(1.5), Inches(1.5), Inches(7), Inches(2.8))
        metric_shape.fill.solid()
        metric_shape.fill.fore_color.rgb = LIGHT_BLUE
        metric_shape.line.color.rgb = DARK_BLUE
        metric_shape.line.width = Pt(4)

        # Metric label
        label_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.7), Inches(7), Inches(0.7))
        label_frame = label_box.text_frame
        p = label_frame.paragraphs[0]
        p.text = metric_title
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Metric value
        value_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(1.3))
        value_frame = value_box.text_frame
        p = value_frame.paragraphs[0]
        p.text = f"{metric_value}\n{metric_unit}"
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description (LARGER for Zoom)
        if description:
            desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(2.5))
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            p = desc_frame.paragraphs[0]
            p.text = description
            p.font.size = Pt(22)
            p.font.color.rgb = DARK_GRAY

        return slide

    def add_comparison_slide(title, left_label, left_value, right_label, right_value,
                            color_left=ACCENT_GREEN, color_right=ACCENT_RED):
        """Add a comparison slide (ZOOM-FRIENDLY)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title box
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = DARK_BLUE
        title_shape.line.color.rgb = DARK_BLUE

        title_frame = title_shape.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Left box
        left_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(4.2), Inches(5.5))
        left_shape.fill.solid()
        left_shape.fill.fore_color.rgb = color_left
        left_shape.line.color.rgb = DARK_GRAY
        left_shape.line.width = Pt(3)

        left_label_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.2), Inches(0.9))
        left_label_frame = left_label_box.text_frame
        left_label_frame.word_wrap = True
        p = left_label_frame.paragraphs[0]
        p.text = left_label
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE

        left_value_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(4.2), Inches(2.5))
        left_value_frame = left_value_box.text_frame
        left_value_frame.word_wrap = True
        p = left_value_frame.paragraphs[0]
        p.text = left_value
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Right box
        right_shape = slide.shapes.add_shape(1, Inches(5.3), Inches(1.5), Inches(4.2), Inches(5.5))
        right_shape.fill.solid()
        right_shape.fill.fore_color.rgb = color_right
        right_shape.line.color.rgb = DARK_GRAY
        right_shape.line.width = Pt(3)

        right_label_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.8), Inches(4.2), Inches(0.9))
        right_label_frame = right_label_box.text_frame
        right_label_frame.word_wrap = True
        p = right_label_frame.paragraphs[0]
        p.text = right_label
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE

        right_value_box = slide.shapes.add_textbox(Inches(5.3), Inches(3.2), Inches(4.2), Inches(2.5))
        right_value_frame = right_value_box.text_frame
        right_value_frame.word_wrap = True
        p = right_value_frame.paragraphs[0]
        p.text = right_value
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        return slide

    # ===== SLIDES START =====

    # SLIDE 1: Title
    add_title_slide(
        "Data Breach Disclosure Timing",
        "& Market Outcomes"
    )

    # SLIDE 2: Central Finding
    add_key_metric_slide(
        "The Core Discovery",
        "Markets Punish:",
        "WHO YOU ARE",
        "& WHAT WAS BREACHED",
        "NOT disclosure timing"
    )

    # SLIDE 3: The Paradox
    add_content_slide(
        "The Timing Paradox",
        [
            "Timing does NOT affect shareholder returns",
            "Timing INCREASES market uncertainty",
            "Timing ACCELERATES governance response",
            "Regulation creates costs without valuation benefits"
        ]
    )

    # SLIDE 4: Dataset
    add_content_slide(
        "The Data",
        [
            "1,054 breaches (2006-2025)",
            "926 firms with stock price data",
            "200 FCC-regulated | 854 non-FCC",
            "442 repeat offenders (42%)"
        ]
    )

    # SLIDE 5: Natural Experiment
    add_content_slide(
        "Research Design: FCC Natural Experiment",
        [
            "FCC Rule 37.3: January 1, 2007",
            "7-day mandate for telecom firms",
            "Control: All other industries",
            "Validation: Effect emerges post-2007, strengthens with controls"
        ]
    )

    # ===== ESSAY 1: MARKET RETURNS =====

    # SLIDE 6: H1 - Timing Effect
    add_key_metric_slide(
        "H1: Does Timing Affect Returns?",
        "Timing Effect",
        "+0.57%",
        "(p=0.539)",
        "NOT significant. Equivalence test confirms genuine null, not power shortage."
    )

    # SLIDE 7: H1 Interpretation
    add_content_slide(
        "What H1 Null Means",
        [
            "Markets are indifferent to disclosure speed",
            "Fast disclosure ≠ lower shareholder losses",
            "Slow disclosure ≠ higher shareholder losses",
            "Refutes core regulatory assumption"
        ]
    )

    # SLIDE 8: H2 - FCC Effect
    add_key_metric_slide(
        "H2: Does FCC Regulation Cost Money?",
        "FCC Regulatory Penalty",
        "-2.20%",
        "CAR (p=0.010)",
        "Significant market penalty for regulatory burden"
    )

    # SLIDE 9: H2 Economic Impact
    add_comparison_slide(
        "FCC Cost: Per Breach vs. Aggregate",
        "Per Breach\n(Median Firm)",
        "-$0.9M",
        "Aggregate\n(187 breaches)",
        "-$0.76B",
        color_left=ACCENT_RED,
        color_right=ACCENT_RED
    )

    # SLIDE 10: H3 - Prior Breaches
    add_key_metric_slide(
        "H3: Prior Breach History (STRONGEST EFFECT)",
        "Per Prior Breach",
        "-0.22%",
        "(p<0.001)",
        "Reputational accumulation. A firm with 5 prior breaches faces -1.1% additional penalty."
    )

    # SLIDE 11: H4 - Health Breaches
    add_key_metric_slide(
        "H4: Does Health Data Cost More?",
        "Health Breach Premium",
        "-2.51%",
        "CAR (p=0.004)",
        "Market-equivalent penalty to FCC regulatory burden"
    )

    # SLIDE 12: Essay 1 Summary
    add_content_slide(
        "Essay 1: What Drives Market Returns?",
        [
            "Prior Breaches: -0.22% per breach ★★★ STRONGEST",
            "FCC Regulation: -2.20% ★★",
            "Health Breach: -2.51% ★★★",
            "Timing: +0.57% (NOT significant)"
        ]
    )

    # ===== ESSAY 2: VOLATILITY =====

    # SLIDE 13: H5 - Volatility Effect
    add_key_metric_slide(
        "H5: Does Mandatory Timing Increase Uncertainty?",
        "Volatility Increase",
        "+1.68%",
        "to +5.02%",
        "Paradoxical: Speed mandate INCREASES uncertainty"
    )

    # SLIDE 14: Volatility Paradox Explained
    add_content_slide(
        "Why Forced Speed Increases Uncertainty",
        [
            "7-day deadline forces disclosure before investigation complete",
            "Incomplete disclosure (\"investigation ongoing\") signals quality problems",
            "Markets see uncertainty, raise volatility",
            "Speed success ≠ Quality success"
        ]
    )

    # SLIDE 15: Size Heterogeneity - Volatility
    add_comparison_slide(
        "Volatility Effect by Firm Size",
        "Small Firms\n(Q1)",
        "+7.31%",
        "Large Firms\n(Q4)",
        "-3.39%",
        color_left=ACCENT_RED,
        color_right=ACCENT_GREEN
    )

    # SLIDE 16: Volatility Interpretation
    add_content_slide(
        "What Size Difference Means",
        [
            "Small firms: Capacity constraint. Cannot investigate AND disclose rapidly.",
            "Large firms: Can accommodate rapid disclosure without quality loss.",
            "Regulatory timing interacts with organizational capacity.",
            "One-size-fits-all rule may harm smaller firms."
        ]
    )

    # ===== ESSAY 3: GOVERNANCE =====

    # SLIDE 17: H6 - Executive Turnover
    add_key_metric_slide(
        "H6: Does Mandatory Disclosure Trigger Governance Changes?",
        "Turnover Acceleration",
        "+5.3 pp",
        "(50.6% vs 45.3%)",
        "Immediate disclosure accelerates executive departures"
    )

    # SLIDE 18: Governance Mechanism
    add_content_slide(
        "The Governance Mechanism",
        [
            "Mandatory disclosure activates stakeholder pressure",
            "Boards respond with executive changes",
            "Turnover serves as accountability signal",
            "Self-response (46%) >> regulatory enforcement (0.6%)"
        ]
    )

    # ===== INTEGRATION =====

    # SLIDE 19: Three Independent Mechanisms
    add_content_slide(
        "Three Mechanisms Operating Independently",
        [
            "Essay 1: Timing irrelevant to valuation",
            "Essay 2: Timing increases uncertainty via quality loss",
            "Essay 3: Timing accelerates governance response",
            "Mediation test proves independence (volatility does NOT mediate governance)"
        ]
    )

    # ===== ECONOMIC SIGNIFICANCE =====

    # SLIDE 20: Cost of Capital
    add_key_metric_slide(
        "Cost of Capital Impact",
        "Volatility-Induced Annual Cost",
        "+$3-4M",
        "per $1B debt",
        "39bp volatility increase for large refinancing"
    )

    # SLIDE 21: Governance Costs
    add_key_metric_slide(
        "Governance Disruption Cost",
        "Per Executive Departure",
        "$12-25M",
        "(total: direct + indirect)",
        "416 breaches with 30-day turnover = $0.39-0.98B aggregate"
    )

    # SLIDE 22: Total Economic Cost
    add_key_metric_slide(
        "Total FCC Regulation Cost",
        "Aggregate Economic Impact",
        "$1.25-1.94B",
        "across three mechanisms",
        "Valuation ($0.76B) + Cost of Capital + Governance"
    )

    # ===== POLICY =====

    # SLIDE 23: FCC Policy Recommendation
    add_content_slide(
        "FCC 7-Day Rule: Policy Implications",
        [
            "Current rule generates $0.76B costs",
            "Timing has no valuation benefit",
            "Forced speed creates quality problems",
            "Recommend: Extend to 14-21 days for investigation"
        ]
    )

    # SLIDE 24: SEC Policy Recommendation
    add_content_slide(
        "SEC 4-Day Rule: Similar Risks",
        [
            "SEC explicitly asked for validation",
            "This dissertation provides evidence",
            "4-day requirement faces identical paradoxes",
            "Recommend: Evaluate 18-24 months post-implementation"
        ]
    )

    # SLIDE 25: Optimal Timeline
    add_content_slide(
        "The Optimal Timeline Concept",
        [
            "Too fast (7 days): Forces incomplete disclosure",
            "Too slow (unregulated): Allows information advantage",
            "Optimal (45-60 days): Allows investigation + stakeholder communication",
            "HIPAA's 60-day window likely superior to FCC's 7-day"
        ]
    )

    # SLIDE 26: Robustness
    add_content_slide(
        "Why These Findings Are Solid",
        [
            "✓ 27+ robustness specifications",
            "✓ Mediation analysis (mechanisms independent)",
            "✓ Temporal validation (pre/post-2007 test)",
            "✓ Machine learning confirms feature importance"
        ]
    )

    # SLIDE 27: Key Takeaways
    add_content_slide(
        "Five Main Findings",
        [
            "1. Timing does NOT reduce shareholder losses",
            "2. Mandatory speed INCREASES market uncertainty",
            "3. Regulation ACCELERATES governance response",
            "4. Economic cost: $0.76B-$1.94B (aggregate)",
            "5. Consider 45-60 day timelines instead of 7 days"
        ]
    )

    # SLIDE 28: Next
    add_content_slide(
        "Next: Dashboard Walkthrough",
        [
            "Interactive filtering by:",
            "• Firm size, FCC status, breach type",
            "• Main results with confidence intervals",
            "• Validation tests (pre/post, industry controls)",
            "• Economic significance calculator"
        ]
    )

    # SLIDE 29: Questions
    add_title_slide(
        "Questions?",
        "Dashboard & Discussion"
    )

    # Save presentation
    output_path = r'C:\Users\mcobp\BA798_TIM\Dissertation_Presentation_Zoom.pptx'
    prs.save(output_path)
    print(f"Zoom-friendly presentation created: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == '__main__':
    create_presentation()
