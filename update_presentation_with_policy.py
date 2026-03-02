"""
Update Dissertation Presentation with Policy Alternatives Slides
Adds 2 slides on the 3 policy frameworks to the existing presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load existing presentation
prs = Presentation(r'C:\Users\mcobp\BA798_TIM\Dissertation_Presentation.pptx')

# Slide 1: Policy Alternatives Overview
slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
title = slide1.shapes.title
title.text = "Policy Alternatives: Addressing Regulatory Costs"
content = slide1.placeholders[1].text_frame
content.text = "Three evidence-based alternatives to current time-based disclosure mandates:\n\n"

# Add bullet points
p = content.add_paragraph()
p.text = "Staged Disclosure Framework"
p.level = 0
p = content.add_paragraph()
p.text = "7-day initial notification + extended investigation window"
p.level = 1
p = content.add_paragraph()
p.text = "Preserves governance activation, allows investigation completeness"
p.level = 1

p = content.add_paragraph()
p.text = "Quality Standards Framework"
p.level = 0
p = content.add_paragraph()
p.text = "Disclosure upon investigation completion OR 7 days, whichever later"
p.level = 1
p = content.add_paragraph()
p.text = "Aligns disclosure timing with completeness rather than arbitrary timelines"
p.level = 1

p = content.add_paragraph()
p.text = "Safe Harbor Framework"
p.level = 0
p = content.add_paragraph()
p.text = "Liability protection for early preliminary disclosures"
p.level = 1
p = content.add_paragraph()
p.text = "Incentivizes early voluntary disclosure, reduces information penalties"
p.level = 1

# Format text
for paragraph in content.paragraphs:
    paragraph.font.size = Pt(16)
    paragraph.font.name = 'Arial'

# Slide 2: Policy Cost-Benefit Analysis
slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
title = slide2.shapes.title
title.text = "Cost-Benefit Analysis of Alternative Frameworks"
content = slide2.placeholders[1].text_frame
content.text = ""

# Staged Framework
p = content.add_paragraph()
p.text = "Staged Disclosure Framework"
p.level = 0
p.font.bold = True
p = content.add_paragraph()
p.text = "Benefits: Governance activation + investigation completeness"
p.level = 1
p = content.add_paragraph()
p.text = "Implementation: Requires coordination with FCC, SEC, HHS"
p.level = 1

p = content.add_paragraph()
p.text = "Quality Standards Framework"
p.level = 0
p.font.bold = True
p = content.add_paragraph()
p.text = "Benefits: Principle-based approach, stakeholder trust enhancement"
p.level = 1
p = content.add_paragraph()
p.text = "Implementation: Simpler regulatory change, market incentive alignment"
p.level = 1

p = content.add_paragraph()
p.text = "Safe Harbor Framework"
p.level = 0
p.font.bold = True
p = content.add_paragraph()
p.text = "Benefits: Incentivizes early voluntary disclosure"
p.level = 1
p = content.add_paragraph()
p.text = "Implementation: Liability law changes, litigation cost redistribution"
p.level = 1

p = content.add_paragraph()
p.text = "Estimated annual shareholder savings across telecommunications sector: $9.9B+"
p.level = 0
p.font.italic = True

# Format text
for paragraph in content.paragraphs:
    paragraph.font.size = Pt(14)
    paragraph.font.name = 'Arial'

# Save updated presentation
prs.save(r'C:\Users\mcobp\BA798_TIM\Dissertation_Presentation_Updated.pptx')
print("[OK] PowerPoint presentation updated with 2 policy alternatives slides")
print("[OK] File: Dissertation_Presentation_Updated.pptx")
print("[OK] Added slides: Policy Alternatives Overview, Cost-Benefit Analysis")
