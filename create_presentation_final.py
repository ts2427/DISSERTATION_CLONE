"""
Final Dissertation Presentation (PowerPoint)
Lean, Zoom-friendly complement to Streamlit dashboard
The PPT is a teaser; the dashboard is the detailed evidence
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_presentation():
    """Generate final lean presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Colors
    DARK_BLUE = RGBColor(31, 78, 121)
    LIGHT_BLUE = RGBColor(79, 129, 189)
    RED = RGBColor(192, 0, 0)
    GREEN = RGBColor(0, 176, 80)
    WHITE = RGBColor(255, 255, 255)
    DARK_GRAY = RGBColor(89, 89, 89)

    def title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = DARK_BLUE

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = WHITE

        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
            p = sub_box.text_frame.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = LIGHT_BLUE

        return slide

    def content_slide(title, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = DARK_BLUE
        title_box.line.color.rgb = DARK_BLUE
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Bullets
        content = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5))
        tf = content.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(24)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(14)
            p.space_after = Pt(14)

        return slide

    def metric_slide(title, metric, value, unit, desc):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = DARK_BLUE
        title_box.line.color.rgb = DARK_BLUE
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Metric box
        metric_box = slide.shapes.add_shape(1, Inches(1.5), Inches(1.5), Inches(7), Inches(2.8))
        metric_box.fill.solid()
        metric_box.fill.fore_color.rgb = LIGHT_BLUE
        metric_box.line.color.rgb = DARK_BLUE
        metric_box.line.width = Pt(4)

        # Metric label
        label_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.7), Inches(7), Inches(0.7))
        p = label_box.text_frame.paragraphs[0]
        p.text = metric
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Value
        value_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(1.3))
        p = value_box.text_frame.paragraphs[0]
        p.text = f"{value}\n{unit}"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(2.5))
        p = desc_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY

        return slide

    # ===== SLIDES =====

    # 1. Title
    title_slide(
        "Data Breach Disclosure Timing",
        "& Market Outcomes"
    )

    # 2. The Question
    content_slide(
        "The Central Question",
        [
            "Is there any benefit to disclosing breaches immediately?",
            "Or should companies delay until investigation is complete?",
            "1,054 breaches, FCC natural experiment, 19 years of data"
        ]
    )

    # 3. The Paradox
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = DARK_BLUE
    title_box.line.color.rgb = DARK_BLUE
    p = title_box.text_frame.paragraphs[0]
    p.text = "The Disclosure Paradox"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content.text_frame
    tf.word_wrap = True

    findings = [
        ("Essay 1:", "Timing does NOT affect shareholder returns", DARK_GRAY),
        ("Essay 2:", "Forced speed INCREASES market uncertainty", RED),
        ("Essay 3:", "Regulation ACCELERATES governance response", GREEN),
    ]

    for i, (label, text, color) in enumerate(findings):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"{label} {text}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*[int(c) for c in color])
        p.space_before = Pt(20)
        p.space_after = Pt(20)

    # 4. Method
    content_slide(
        "Why We Trust These Results",
        [
            "Natural experiment: FCC Rule 37.3 (January 2007)",
            "7-day mandate for telecom → treatment group",
            "All other industries → control group",
            "Validation: Effect emerges post-2007, strengthens with controls"
        ]
    )

    # 5. Essay 1 Finding
    metric_slide(
        "Essay 1: Market Returns",
        "What Markets Actually Care About",
        "Firm Characteristics",
        "NOT timing speed",
        "Prior breach history (-0.22%/breach) and regulatory status (-2.20%) drive returns, not timing"
    )

    # 6. Essay 2 Finding
    metric_slide(
        "Essay 2: Information Asymmetry",
        "Mandatory Speed Creates",
        "HIGHER Volatility",
        "+1.68% to +5.02%",
        "Forced speed prevents investigation completion. Incomplete disclosures INCREASE uncertainty."
    )

    # 7. Essay 3 Finding
    metric_slide(
        "Essay 3: Governance Response",
        "Firms Respond With",
        "Executive Turnover",
        "+5.3 percentage points",
        "46.4% of breaches trigger exec departures within 30 days via stakeholder pressure"
    )

    # 8. Economic Impact
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = DARK_BLUE
    title_box.line.color.rgb = DARK_BLUE
    p = title_box.text_frame.paragraphs[0]
    p.text = "Economic Significance"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Two boxes
    left = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(4.2), Inches(5.5))
    left.fill.solid()
    left.fill.fore_color.rgb = RED
    left.line.color.rgb = DARK_GRAY
    left.line.width = Pt(3)

    p = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.2), Inches(0.8)).text_frame.paragraphs[0]
    p.text = "Per Breach\n(Median Firm)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(4.2), Inches(2)).text_frame.paragraphs[0]
    p.text = "-$0.9M"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    right = slide.shapes.add_shape(1, Inches(5.3), Inches(1.5), Inches(4.2), Inches(5.5))
    right.fill.solid()
    right.fill.fore_color.rgb = RED
    right.line.color.rgb = DARK_GRAY
    right.line.width = Pt(3)

    p = slide.shapes.add_textbox(Inches(5.3), Inches(1.8), Inches(4.2), Inches(0.8)).text_frame.paragraphs[0]
    p.text = "Aggregate\n(187 FCC breaches)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p = slide.shapes.add_textbox(Inches(5.3), Inches(3), Inches(4.2), Inches(2)).text_frame.paragraphs[0]
    p.text = "-$0.76B"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 9. Policy for FCC
    content_slide(
        "FCC 7-Day Rule: Key Findings",
        [
            "✓ Achieves goal: Disclosure happens faster",
            "✗ Creates problem: Forced speed increases uncertainty",
            "✗ No benefit: Timing does not reduce shareholder losses",
            "Recommendation: Consider 14-21 day window (investigate fully)"
        ]
    )

    # 10. Policy for SEC
    content_slide(
        "SEC 4-Day Rule: Similar Story",
        [
            "SEC asked for empirical validation of 4-day timeline",
            "This research shows it faces identical paradoxes",
            "Speed ≠ Information quality",
            "Recommendation: Monitor implementation impact closely"
        ]
    )

    # 11. Business Implications
    content_slide(
        "For Breached Companies",
        [
            "Early disclosure DOES accelerate governance response",
            "Early disclosure DOES NOT reduce market losses",
            "Timing is governance signal, not valuation signal",
            "Focus on COMPLETENESS, not just SPEED"
        ]
    )

    # 12. What We DON'T Cover
    content_slide(
        "Important Scope Notes",
        [
            "✓ Evidence = shareholder returns only",
            "✗ Consumer protection outcomes beyond scope",
            "✗ Employee morale beyond scope",
            "✗ Regulatory compliance quality beyond scope"
        ]
    )

    # 13. Robustness
    content_slide(
        "Why These Findings Are Solid",
        [
            "✓ 27+ robustness specifications (all support findings)",
            "✓ Temporal validation: Effect emerges exactly post-2007",
            "✓ Mediation tests: Three mechanisms operate independently",
            "✓ Machine learning: Feature importance confirms results"
        ]
    )

    # 14. To Explore Details
    content_slide(
        "Let's Explore the Dashboard",
        [
            "📊 15-page interactive Streamlit dashboard",
            "🔍 Filter by: FCC status, firm size, breach type",
            "📈 Detailed regression results with confidence intervals",
            "📋 All validation tests & robustness checks visible"
        ]
    )

    # 15. Dashboard Structure
    content_slide(
        "Dashboard Navigation",
        [
            "Pages 1-3: Context & Natural Experiment",
            "Pages 4-6: Essay 1-3 detailed results",
            "Pages 7-8: Economic significance & Robustness",
            "Pages 9-11: Conclusions, Data Explorer, Dictionary"
        ]
    )

    # 16. Questions
    title_slide(
        "Questions?",
        "Exploring the Dashboard Together"
    )

    # Save
    output_path = r'C:\Users\mcobp\BA798_TIM\Dissertation_Presentation.pptx'
    prs.save(output_path)
    print(f"Final presentation: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == '__main__':
    create_presentation()
