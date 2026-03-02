"""
Expanded Dissertation Presentation (28 slides)
Comprehensive but not gratuitous - covers all major content from proposal
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_presentation():
    """Generate expanded presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Colors
    DARK_BLUE = RGBColor(31, 78, 121)
    LIGHT_BLUE = RGBColor(79, 129, 189)
    RED = RGBColor(192, 0, 0)
    GREEN = RGBColor(0, 176, 80)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(89, 89, 89)

    def title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = DARK_BLUE

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = WHITE

        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
            p = sub_box.text_frame.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = LIGHT_BLUE

        return slide

    def content_slide(title, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = DARK_BLUE
        title_box.line.color.rgb = DARK_BLUE
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE

        content = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5))
        tf = content.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(22)
            p.font.color.rgb = GRAY
            p.space_before = Pt(12)
            p.space_after = Pt(12)

        return slide

    def metric_slide(title, metric, value, unit, desc=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = DARK_BLUE
        title_box.line.color.rgb = DARK_BLUE
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE

        metric_box = slide.shapes.add_shape(1, Inches(1.5), Inches(1.5), Inches(7), Inches(2.8))
        metric_box.fill.solid()
        metric_box.fill.fore_color.rgb = LIGHT_BLUE
        metric_box.line.color.rgb = DARK_BLUE
        metric_box.line.width = Pt(4)

        label_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.7), Inches(7), Inches(0.7))
        p = label_box.text_frame.paragraphs[0]
        p.text = metric
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE

        value_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(1.3))
        p = value_box.text_frame.paragraphs[0]
        p.text = f"{value}\n{unit}"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        if desc:
            desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(2.5))
            p = desc_box.text_frame.paragraphs[0]
            p.text = desc
            p.font.size = Pt(19)
            p.font.color.rgb = GRAY

        return slide

    # ===== SECTION 1: INTRODUCTION & CONTEXT =====

    # 1. Title
    title_slide(
        "Data Breach Disclosure Timing",
        "& Market Outcomes"
    )

    # 2. The Problem
    content_slide(
        "The Problem: Regulatory Assumptions Without Evidence",
        [
            "FCC mandates 7-day breach disclosure (Rule 37.3, 2007)",
            "SEC mandates 4-day disclosure (2023 rule)",
            "HIPAA mandates 60-day disclosure",
            "All based on assumption: Faster = Better",
            "But: NO causal evidence exists that speed creates benefits"
        ]
    )

    # 3. Why It Matters
    content_slide(
        "Why This Research Matters",
        [
            "For Regulators: FCC rule generates $0.76B in shareholder losses",
            "For Companies: Disclosure timing is a strategic decision affecting shareholders, employees, customers",
            "For Theory: Information asymmetry theory predicts timing effects, but does it hold?",
            "For Markets: Do investors reward speed or penalize incomplete disclosures?"
        ]
    )

    # 4. Literature Foundation
    content_slide(
        "What Prior Research Shows",
        [
            "Event Studies: Breaches cause −0.41% to −2.1% returns (Cavusoglu 2004, Acquisti 2006)",
            "Paradox: Mandatory disclosure laws increase crash risk 5−7% (Obaydin et al. 2024)",
            "Theory Gap: Diamond & Verrecchia (1991) warns forced disclosure can INCREASE asymmetry",
            "Missing: Clean causal evidence separating timing effects from firm quality signals"
        ]
    )

    # ===== SECTION 2: RESEARCH DESIGN =====

    # 5. Sample Overview
    content_slide(
        "The Dataset: 1,054 Breaches (2006-2025)",
        [
            "926 breaches with stock price data (Essay 1)",
            "916 breaches with volatility data (Essay 2)",
            "896 breaches with executive change data (Essay 3)",
            "200 FCC-regulated firms (19%) | 854 non-FCC (81%)",
            "442 repeat offenders (41.9% with prior breach history)"
        ]
    )

    # 6. Natural Experiment
    content_slide(
        "Causal Identification: FCC Rule 37.3 Natural Experiment",
        [
            "FCC Rule effective January 1, 2007: 7-day mandate for telecom (SIC 4813, 4899, 4841)",
            "Control: All other industries under state-level rules (30−90 day timelines)",
            "Sharp regulatory discontinuity creates quasi-experimental variation",
            "Can compare FCC vs. non-FCC breach outcomes before/after 2007"
        ]
    )

    # 7. Validation Tests
    content_slide(
        "Three Tests Validate Causal Interpretation",
        [
            "Test 1 (Temporal): FCC effect zero pre-2007 (p=0.88), emerges post-2007 (p=0.0125) ✓",
            "Test 2 (Controls): Effect strengthens with industry FE (−2.20% → −5.37%), not selection ✓",
            "Test 3 (Size Heterogeneity): FCC effects concentrated in small firms (−6.22%), null in large ✓",
            "Multi-outcome consistency: FCC affects returns, volatility, AND governance"
        ]
    )

    # ===== SECTION 3: HYPOTHESES =====

    # 8. H1 - Timing Effect
    metric_slide(
        "H1: Does Timing Affect Shareholder Returns?",
        "Timing Coefficient",
        "+0.57%",
        "(p=0.539)",
        "NOT significant. Markets are indifferent to disclosure speed. Equivalence test confirms genuine null."
    )

    # 9. H1 Interpretation
    content_slide(
        "H1 Null: What It Means",
        [
            "Disclosure speed does NOT predict shareholder losses",
            "Fast disclosure ≠ lower shareholder penalty",
            "Slow disclosure ≠ higher shareholder penalty",
            "Refutes core regulatory assumption",
            "Markets efficiently price breach fundamentals regardless of timing"
        ]
    )

    # 10. H2 - FCC Regulation
    metric_slide(
        "H2: Does FCC Regulation Cost Money?",
        "FCC Regulatory Penalty",
        "-2.20%",
        "CAR (p=0.010)",
        "Significant market penalty. FCC-regulated firms experience ~$4.0M average loss per breach."
    )

    # 11. H3 - Prior Breaches (Strongest)
    metric_slide(
        "H3: Reputation Accumulation (STRONGEST EFFECT)",
        "Per Prior Breach",
        "-0.22%",
        "(p<0.001)",
        "A firm with 5 prior breaches faces −1.1% additional penalty on the next breach."
    )

    # 12. H4 - Health Breaches
    metric_slide(
        "H4: Does Health Data Cost More?",
        "Health Breach Premium",
        "-2.51%",
        "CAR (p=0.004)",
        "Health breaches produce nearly identical market penalty to FCC regulatory burden."
    )

    # 13. H5 - Volatility Paradox
    metric_slide(
        "H5: Does Mandatory Timing Increase Uncertainty?",
        "Volatility Effect",
        "+1.68%",
        "to +5.02%",
        "Paradoxical: Forced speed requirement INCREASES market uncertainty instead of decreasing it."
    )

    # 14. H5 Mechanism
    content_slide(
        "Why Forced Speed Increases Uncertainty",
        [
            "7-day deadline forces disclosure before investigation completion",
            "Incomplete disclosure (\"investigation ongoing\") signals quality problems to markets",
            "Markets respond with higher volatility, not lower",
            "Small firms most affected (+7.31%***); large firms absorb speed (−3.39%**)",
            "Regulatory timing FAILS at information quality goal"
        ]
    )

    # 15. H6 - Executive Turnover
    metric_slide(
        "H6: Does Mandatory Disclosure Trigger Governance Changes?",
        "Turnover Acceleration",
        "+5.3pp",
        "(50.6% vs 45.3%)",
        "Immediate disclosure accelerates executive departures via stakeholder pressure activation."
    )

    # 16. H6 Mechanism
    content_slide(
        "Governance Response Mechanism",
        [
            "Mandatory disclosure makes regulator a \"definitive\" stakeholder",
            "Boards respond with executive changes as accountability signal",
            "46.4% of breaches experience 30-day turnover (416 breaches)",
            "Governance self-response >> regulatory enforcement (50:1 ratio, 46% vs 0.6%)",
            "This operates independently of information content"
        ]
    )

    # ===== SECTION 4: INTEGRATION & ROBUSTNESS =====

    # 17. Three Independent Mechanisms
    content_slide(
        "Three Mechanisms Operate Independently",
        [
            "Essay 1 (Valuation): Timing irrelevant → markets price firm characteristics",
            "Essay 2 (Uncertainty): Timing increases volatility → forced speed reduces quality",
            "Essay 3 (Governance): Timing accelerates turnover → stakeholder pressure, not info quality",
            "Mediation test: Volatility does NOT mediate governance (p=0.484)",
            "Essays 2 & 3 are completely separate channels"
        ]
    )

    # 18. Heterogeneous Effects
    content_slide(
        "Size Heterogeneity: Reveals Mechanism",
        [
            "Market Returns (Essay 1): FCC penalty concentrated in small firms (Q1: −6.22%)",
            "Volatility (Essay 2): OPPOSITE pattern - small firms +7.31%, large firms −3.39%",
            "Governance (Essay 3): U-shaped by size (medium firms most responsive)",
            "Interpretation: Information processing capacity constraints. Small firms can't investigate rapidly AND disclose completely."
        ]
    )

    # 19. Validation: Robustness
    content_slide(
        "27+ Robustness Specifications Support Findings",
        [
            "✓ 4 event windows (1d, 5d, 10d, 30d) - consistent patterns",
            "✓ 7 timing thresholds - H1 null holds across all definitions",
            "✓ 8 subsamples (FCC, non-FCC, health, financial, etc.) - effects robust",
            "✓ 6 SE methods (HC3, clustered, double-clustered) - not sensitive to specification",
            "✓ Machine learning validation: Prior breach history #1 feature importance, timing much lower"
        ]
    )

    # ===== SECTION 5: ECONOMIC SIGNIFICANCE & POLICY =====

    # 20. Economic Impact
    content_slide(
        "Economic Significance: FCC Regulation Cost",
        [
            "Per breach (median firm): −$0.9M shareholder loss",
            "Aggregate (187 FCC-regulated breaches): −$0.76 BILLION",
            "Cost of capital increase: +39bp volatility = +$3−4M annually per $1B debt",
            "Governance disruption: $12−25M per executive departure",
            "Total: $1.25−1.94B across all mechanisms"
        ]
    )

    # 21. FCC Policy
    content_slide(
        "FCC 7-Day Rule: Policy Recommendation",
        [
            "Current rule achieves SPEED goal: Forces faster disclosure",
            "But FAILS QUALITY goal: Forced speed prevents thorough investigation",
            "Market consequence: $0.76B shareholder losses without valuation benefits",
            "Recommendation: Extend to 14−21 days to allow investigation completion",
            "Alternative: Conditional disclosure (\"preliminary findings pending investigation\")"
        ]
    )

    # 22. SEC Policy
    content_slide(
        "SEC 4-Day Cybersecurity Rule (2023): Identical Risks",
        [
            "SEC explicitly requested empirical validation before implementation",
            "This dissertation provides that evidence",
            "SEC rule faces same paradoxes as FCC: speed vs. quality tradeoff",
            "Recommendation: Implement with built-in 18−24 month evaluation",
            "Budget IT systems for rapid investigation capability"
        ]
    )

    # 23. Optimal Timeline Concept
    content_slide(
        "The Optimal Disclosure Timeline",
        [
            "Too Fast (7 days): Forces incomplete disclosure → increases uncertainty",
            "Too Slow (unregulated): Allows information asymmetry advantage",
            "Optimal (45−60 days): Allows investigation completion + stakeholder communication",
            "Evidence: HIPAA's 60-day window likely superior to FCC's 7-day",
            "Stakeholder research (Xu et al. 2024): Prefer completeness over speed"
        ]
    )

    # 24. Business Implications
    content_slide(
        "For Breached Companies",
        [
            "Early disclosure accelerates governance response (H6 effect)",
            "Early disclosure does NOT reduce shareholder market losses (H1 null)",
            "Optimal strategy: Disclose for governance credibility, not valuation protection",
            "Focus on COMPLETENESS of disclosure, not raw speed",
            "Timing is governance signal to stakeholders, not market signal to investors"
        ]
    )

    # ===== SECTION 6: LIMITATIONS & CONCLUSIONS =====

    # 25. Limitations
    content_slide(
        "Limitations & Scope Boundaries",
        [
            "Sample: FCC firms 2x larger; size controls address this, but parallel trends may fail",
            "Scope: Evidence is shareholder returns ONLY; consumer protection beyond scope",
            "Timing windows: 30/90/180-day governance windows may miss longer-run effects",
            "Public firms only: Results don't generalize to private company dynamics",
            "Despite limitations: Three-essay design + natural experiment + 27+ robustness tests"
        ]
    )

    # 26. Key Takeaways
    content_slide(
        "Five Main Conclusions",
        [
            "1. MARKETS: Disclosure timing irrelevant; firm characteristics drive returns",
            "2. UNCERTAINTY: Mandatory speed increases volatility through quality tradeoff",
            "3. GOVERNANCE: Regulation accelerates governance response via stakeholder pressure",
            "4. POLICY: Current 7-day rule creates costs without valuation benefits",
            "5. THEORY: Information asymmetry requires quality, not just speed"
        ]
    )

    # 27. Dashboard
    content_slide(
        "Interactive Dashboard: Deep Dive",
        [
            "15-page Streamlit application with all detailed evidence",
            "Filter by: FCC status, firm size, breach type, disclosure timing",
            "Pages 1−3: Context, Natural Experiment, Sample Validation",
            "Pages 4−6: Essay results with tables and confidence intervals",
            "Pages 7−11: Robustness, heterogeneous effects, conclusions, data explorer"
        ]
    )

    # 28. Questions
    title_slide(
        "Questions?",
        "Exploring the Dashboard Together"
    )

    # Save
    output_path = r'C:\Users\mcobp\BA798_TIM\Dissertation_Presentation.pptx'
    prs.save(output_path)
    print(f"Expanded presentation: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == '__main__':
    create_presentation()
