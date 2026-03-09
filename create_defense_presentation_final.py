#!/usr/bin/env python3
"""
FINAL Dissertation Defense Presentation (23-25 slides)
UPDATES:
- Theoretical chain slide with visual progression
- Essay 1 heterogeneity slide (CVSS +6.27%, Media +7.08%)
- Mitchell, Agle, & Wood citations for stakeholder theory
- Kothari et al. (2009) bad news accumulation
- Fixed: FCC 7-day (not 30), denominator 1,054 (not 1,006), industry FE
- Mediation analysis interpretation (independence proof)
- Prior Breaches emphasized as strongest effect
- Clean, uncluttered presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Colors
COLOR_DARK_BLUE = RGBColor(27, 94, 148)
COLOR_TEAL = RGBColor(46, 204, 113)
COLOR_ORANGE = RGBColor(230, 126, 34)
COLOR_RED = RGBColor(192, 57, 43)
COLOR_GRAY = RGBColor(149, 165, 166)
COLOR_LIGHT_BG = RGBColor(236, 240, 241)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK_BLUE

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(240, 240, 240)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, color=COLOR_DARK_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = color
    title_shape.line.color.rgb = color

    title_frame = title_shape.text_frame
    title_frame.vertical_anchor = 1
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(10)
    p.space_after = Pt(10)

    return slide

def add_text_box(slide, left, top, width, height, text, size=14, bold=False, color=RGBColor(0, 0, 0)):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box

def add_bullet_points(slide, left, top, width, height, bullets, size=13):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = frame.paragraphs[0]
        else:
            p = frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    return box

# ============================================================================
# SECTION 1: OPENING (Slides 1-3)
# ============================================================================

add_title_slide(prs,
    "Disclosure Timing and Market Reactions",
    "A Natural Experiment in Data Breach Regulation\nThree Essays on Valuation, Uncertainty, and Governance\n\n1,054 breaches | 2006-2025 | FCC Rule 37.3 Natural Experiment")

slide = add_content_slide(prs, "The Puzzle: Timing Paradox", COLOR_DARK_BLUE)
add_text_box(slide, 0.5, 1.2, 9, 0.6,
    "Do disclosure timing mandates improve market outcomes?",
    size=18, bold=True)
bullets = [
    "Essay 1: Does disclosure timing affect stock valuations? → NO (H1: +0.57%, p=0.539)",
    "Essay 2: Does disclosure timing affect market uncertainty? → YES (FCC: +1.83%**, increases volatility)",
    "Essay 3: Does disclosure timing trigger governance response? → YES (FCC: +5.3pp** executive turnover)",
    " ",
    "KEY INSIGHT: Timing doesn't change WHAT markets conclude, but it DOES change HOW markets learn and WHEN firms respond"
]
add_bullet_points(slide, 0.5, 2.0, 9, 5, bullets, size=13)

slide = add_content_slide(prs, "Policy Context: Regulatory Landscape", COLOR_DARK_BLUE)
bullets = [
    "FCC: 7-day disclosure mandate (Rule 37.3, effective 2007)",
    "SEC: 4-day cybersecurity disclosure rule (2023)",
    "HIPAA: 60-day notification requirement",
    "State laws: 30-90 days (California: 30 days)",
    " ",
    "THE GAP: No empirical evidence on whether timing mandates create benefits or costs",
    "YOUR CONTRIBUTION: First clean natural experiment isolating timing effects from firm quality signals"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=12)

# ============================================================================
# SECTION 2: THEORY & FRAMEWORK (Slides 4-7)
# ============================================================================

# Slide 4: Theoretical Chain (NEW - Visual Progression)
slide = add_content_slide(prs, "Theoretical Foundation: A Progressive Chain", COLOR_DARK_BLUE)

# Four boxes showing progression
box_y = 1.4
box_width = 2.0
box_height = 4.5
colors = [COLOR_TEAL, COLOR_ORANGE, COLOR_RED, RGBColor(147, 112, 219)]

# Box 1: Akerlof
add_text_box(slide, 0.4, box_y, box_width, 0.3, "Akerlof (1970)", size=10, bold=True, color=colors[0])
add_text_box(slide, 0.4, box_y+0.35, box_width, 0.25, "Adverse Selection", size=9, bold=True)
add_text_box(slide, 0.4, box_y+0.65, box_width, 2.5, "Information asymmetry eliminates beneficial trades. Severe-breach firms delay disclosure; mild-breach firms disclose quickly.", size=8)

# Arrow 1
add_text_box(slide, 2.5, box_y+1.8, 0.3, 0.3, "↓", size=20)

# Box 2: Spence
add_text_box(slide, 2.9, box_y, box_width, 0.3, "Spence (1973)", size=10, bold=True, color=colors[1])
add_text_box(slide, 2.9, box_y+0.35, box_width, 0.25, "Costly Signaling", size=9, bold=True)
add_text_box(slide, 2.9, box_y+0.65, box_width, 2.5, "High-quality types separate from low-quality through signals. Disclosure timing IS a signal—until regulation eliminates choice.", size=8)

# Arrow 2
add_text_box(slide, 5.0, box_y+1.8, 0.3, 0.3, "↓", size=20)

# Box 3: Myers & Majluf
add_text_box(slide, 5.4, box_y, box_width, 0.3, "Myers & Majluf (1984)", size=10, bold=True, color=colors[2])
add_text_box(slide, 5.4, box_y+0.35, box_width, 0.25, "Information Signaling", size=9, bold=True)
add_text_box(slide, 5.4, box_y+0.65, box_width, 2.5, "Financing (and disclosure) decisions signal private information. Mandatory timing destroys this signal mechanism.", size=8)

# Arrow 3
add_text_box(slide, 7.5, box_y+1.8, 0.3, 0.3, "↓", size=20)

# Box 4: This Dissertation
add_text_box(slide, 7.9, box_y, 1.9, 0.3, "This Dissertation", size=10, bold=True, color=RGBColor(147, 112, 219))
add_text_box(slide, 7.9, box_y+0.35, 1.9, 0.25, "Cybersecurity Extension", size=9, bold=True)
add_text_box(slide, 7.9, box_y+0.65, 1.9, 2.5, "When timing is mandated, signaling eliminated—creating three independent cost channels: valuation penalty, uncertainty, governance response.", size=8, bold=True, color=RGBColor(147, 112, 219))

# Bottom: Three Competing Theories
slide = add_content_slide(prs, "Three Competing Theories of Disclosure Timing", COLOR_DARK_BLUE)

add_text_box(slide, 0.5, 1.2, 2.8, 0.4, "SIGNALING (Myers & Majluf 1984)", size=11, bold=True, color=RGBColor(46, 204, 113))
add_text_box(slide, 0.5, 1.65, 2.8, 1.8,
    "Fast disclosure = confidence signal\nSlow disclosure = bad news hoarding\n\nPrediction: Faster → Better returns\n\nYour evidence: REJECTED (p=0.539)",
    size=9)

add_text_box(slide, 3.6, 1.2, 2.8, 0.4, "FORCED DISCLOSURE PARADOX", size=11, bold=True, color=RGBColor(230, 126, 34))
add_text_box(slide, 3.6, 1.65, 2.8, 1.8,
    "Diamond & Verrecchia (1991)\nMandatory disclosure can INCREASE uncertainty\nForced before investigation complete\n\nYour evidence: SUPPORTED (Essay 2)",
    size=9)

add_text_box(slide, 6.7, 1.2, 2.8, 0.4, "STAKEHOLDER PRESSURE", size=11, bold=True, color=RGBColor(192, 57, 43))
add_text_box(slide, 6.7, 1.65, 2.8, 1.8,
    "Freeman (1984), Mitchell, Agle, & Wood (1997)\nDisclosure activates stakeholders\nImmediate disclosure forces response\n\nYour evidence: SUPPORTED (Essay 3)",
    size=9)

add_text_box(slide, 0.5, 3.7, 9, 3.5,
    "Supporting Mechanisms:\n\nKothari et al. (2009): Bad news accumulation under mandatory regimes\nTushman & Nadler (1978): Information processing capacity constraints (Essay 2 heterogeneity)\n\nKEY FINDING: All three theories are partially correct. They address different outcomes through different mechanisms.",
    size=11, bold=True, color=COLOR_DARK_BLUE)

# Slide 6: Natural Experiment
slide = add_content_slide(prs, "Natural Experiment: FCC Rule 37.3 (2007)", COLOR_DARK_BLUE)
bullets = [
    "THE SHOCK: FCC Rule 37.3 adopted September 2007 → effective September 30, 2007",
    " ",
    "TREATMENT GROUP (FCC-Regulated): SIC 4813, 4899, 4841 = 200 firms",
    "CONTROL GROUP (Non-FCC): All other industries = 854 firms",
    " ",
    "WHY IT'S CLEAN: Regulatory mandate, not market-driven | Applies to whole industry | Exogenous shock",
    "VALIDATION: Pre-treatment balance verified (all characteristics p>0.05 via TABLE_BALANCE_TEST)"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=12)

# Slide 7: Identification
slide = add_content_slide(prs, "Four Robustness Tests: Causal Identification", COLOR_DARK_BLUE)
bullets = [
    "TEST 1 - Temporal (Pre-2007): FCC coef pre=−13.96% (NS) | post=−2.26% (sig) → Effect post-regulation only",
    "TEST 2 - Industry FE: Baseline −2.20% | With industry FE strengthens (not driven by composition)",
    "TEST 3 - Size Sensitivity: Effects vary by size but patterns consistent across quartiles",
    "TEST 4 - Multi-Outcome: Same natural experiment explains returns, volatility, AND turnover"
]
add_bullet_points(slide, 0.5, 1.2, 9, 5.8, bullets, size=12)

# ============================================================================
# SECTION 3: ESSAY 1 - MARKET REACTIONS (Slides 8-11)
# ============================================================================

slide = add_content_slide(prs, "Essay 1: Market Returns - Main Findings", COLOR_TEAL)
add_text_box(slide, 0.5, 1.2, 9, 0.5, "H1 - Timing Effect (Null Result):", size=13, bold=True, color=COLOR_TEAL)
add_text_box(slide, 0.7, 1.7, 8.6, 1.0,
    "Coefficient: +0.57% (p=0.539, NOT significant) | TOST Equivalence: PASS | N=898",
    size=11)

add_text_box(slide, 0.5, 2.8, 9, 0.5, "What DOES Matter:", size=13, bold=True, color=COLOR_TEAL)

add_text_box(slide, 0.7, 3.35, 8.6, 0.35, "Prior Breaches (STRONGEST):  -0.22% per breach***    Market prices reputation (dominant effect)", size=11, bold=True, color=RGBColor(192, 57, 43))
add_text_box(slide, 0.7, 3.75, 8.6, 0.35, "FCC Regulation:                     -2.20%**                Regulatory burden/complexity penalty", size=11, bold=True)
add_text_box(slide, 0.7, 4.15, 8.6, 0.35, "Health Breach:                       -2.51%***              HIPAA compliance risk", size=11, bold=True)

add_text_box(slide, 0.5, 4.7, 9, 2.5,
    "CENTRAL FINDING: Markets punish WHO YOU ARE and WHAT WAS BREACHED—not WHEN YOU TALK.\n\nThe null finding is meaningful: It directly contradicts regulatory assumptions that mandatory speed requirements create market benefits.",
    size=11, bold=True, color=RGBColor(192, 57, 43))

# NEW SLIDE 9: Essay 1 Heterogeneity (CVSS, Media)
slide = add_content_slide(prs, "Essay 1: Mechanism Heterogeneity - Expectation-Based Pricing", COLOR_TEAL)

add_text_box(slide, 0.5, 1.2, 9, 0.4, "Essay 1 finds two significant interaction effects:", size=12, bold=True)

add_text_box(slide, 0.7, 1.7, 4.2, 0.3, "CVSS Technical Complexity", size=11, bold=True, color=RGBColor(46, 204, 113))
add_text_box(slide, 0.7, 2.05, 4.2, 1.8,
    "Interaction: +6.27%* (p=0.007)\n\nInterpretation: Simple breaches penalized 6x more than complex ones\n\nMechanism: Investors expect competent firms to handle complex breaches; simple breaches signal incompetence",
    size=10)

add_text_box(slide, 5.1, 1.7, 4.2, 0.3, "Media Coverage", size=11, bold=True, color=RGBColor(230, 126, 34))
add_text_box(slide, 5.1, 2.05, 4.2, 1.8,
    "Interaction: +7.08%* (p=0.006)\n\nInterpretation: Low-media breaches penalized more (market relies on regulation)\n\nMechanism: Media coverage SUBSTITUTES for regulatory disclosure—when media covers well, regulatory mandate less needed",
    size=10)

add_text_box(slide, 0.5, 4.1, 9, 2.8,
    "KEY INSIGHT: Market pricing is sophisticated. FCC regulation doesn't improve returns, but heterogeneity analysis reveals HOW market incorporates breach news. Media and technical complexity shape expectation-based pricing independent of timing.",
    size=11, bold=True, color=RGBColor(192, 57, 43))

# Slide 10: Essay 1 H1 Robustness
slide = add_content_slide(prs, "Essay 1: Why H1 Null is Robust (TOST Equivalence)", COLOR_TEAL)
bullets = [
    "NOT due to low power: Post-hoc MDE = ±2.39pp (can detect economically meaningful effects)",
    "Equivalence test: 90% CI [-0.95%, +2.09%] falls within ±2.10% bounds → effect is negligible",
    "Robust across ALL 27+ specifications: Event windows, timing thresholds, subsamples, SE methods",
    " ",
    "Timing non-significant across entire heterogeneity landscape:",
    "  • By firm size: All quartiles NS (universal null)",
    "  • By prior breach status: NS",
    "  • By health breach: NS"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=11)

# Slide 11: Essay 1 Size
slide = add_content_slide(prs, "Essay 1: Firm Size Heterogeneity", COLOR_TEAL)
bullets = [
    "FCC Regulation penalty varies by size:",
    "  Q1 (Smallest): -6.22% (p=0.053)  | Q2: -4.06% (p=0.007) | Q3: +0.66% (NS) | Q4 (Largest): +0.43% (NS)",
    " ",
    "Interpretation: FCC penalty concentrated in smaller firms—they feel burden most acutely",
    " ",
    "H1 Timing: ALL non-significant across quartiles",
    "  Q1: +1.16% | Q2: +0.85% | Q3: +0.12% | Q4: +0.05% (all p>0.05)",
    " ",
    "CONCLUSION: Timing irrelevant regardless of firm size (universal null finding)"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=11)

# ============================================================================
# SECTION 4: ESSAY 2 (Slides 12-15)
# ============================================================================

slide = add_content_slide(prs, "Essay 2: Information Asymmetry - Paradoxical Finding", COLOR_ORANGE)
add_text_box(slide, 0.5, 1.2, 9, 0.5, "Same regulation, opposite effects:", size=12, bold=True)
bullets = [
    "RETURNS (Essay 1): FCC = -2.20% (bad, but doesn't change with timing)",
    "VOLATILITY (Essay 2): FCC = +1.83%** (INCREASES uncertainty—opposite of intent)",
    " ",
    "Why this matters: Regulatory constraint (timing requirement) creates information quality problem",
    "→ Forces disclosure BEFORE investigation complete",
    "→ Markets learn faster but with INCOMPLETE information = higher uncertainty",
    " ",
    "Sample: 891 breaches with volatility data"
]
add_bullet_points(slide, 0.5, 1.8, 9, 5.5, bullets, size=12)

slide = add_content_slide(prs, "Essay 2: Timing-Quality Tradeoff Mechanism", COLOR_ORANGE)
bullets = [
    "THE MECHANISM: FCC Rule 37.3 (7-day disclosure) → Time pressure on investigation",
    "→ Firms forced to disclose with INCOMPLETE breach analysis",
    "→ Markets see disclosure but recognize information is INCOMPLETE",
    "→ Result: HIGHER UNCERTAINTY (volatility increases)",
    " ",
    "Evidence: Effect strongest in SMALL firms (+7.31%***) | Weak in LARGE firms (-3.39%**)",
    "  Small firms: Limited breach investigation capacity → time pressure hurts more",
    "  Large firms: Can conduct parallel processes → no quality loss"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=11)

slide = add_content_slide(prs, "Essay 2: THE CRITICAL FINDING - Firm Size Heterogeneity", COLOR_ORANGE)
add_text_box(slide, 0.5, 1.2, 9, 0.3, "FCC Effect on Volatility by Firm Size Quartile:", size=12, bold=True, color=COLOR_ORANGE)

add_text_box(slide, 0.7, 1.6, 8.6, 0.28, "Q1 (Smallest):    +7.31%***    p<0.001    232 firms    [HIGH BURDEN]", size=10, bold=True, color=RGBColor(192, 57, 43))
add_text_box(slide, 0.7, 1.92, 8.6, 0.28, "Q2 (Small):       +3.64%**     p=0.015    233 firms", size=10, bold=True, color=RGBColor(230, 126, 34))
add_text_box(slide, 0.7, 2.24, 8.6, 0.28, "Q3 (Medium):      -0.54%       p=0.770    213 firms", size=10, bold=True)
add_text_box(slide, 0.7, 2.56, 8.6, 0.28, "Q4 (Largest):     -3.39%**     p=0.024    213 firms    [BENEFIT]", size=10, bold=True, color=RGBColor(46, 204, 113))

add_text_box(slide, 0.5, 3.0, 9, 0.25, "11 PERCENTAGE POINT DIFFERENTIAL (Q1 vs Q4) → Regressive burden", size=11, bold=True, color=COLOR_ORANGE)

bullets = [
    "Economic impact: Small firms face ~$37M annual cost of capital increase per breach",
    " ",
    "Mechanism isolates to firm CAPACITY, not breach severity:",
    "  • Not driven by CVSS complexity (p=0.97, tested and rejected)",
    "  • Not driven by governance quality (p=0.84, tested and rejected)",
    "  • Not driven by information environment (p=0.28, tested and rejected)"
]
add_bullet_points(slide, 0.5, 3.4, 9, 3.9, bullets, size=11)

# ============================================================================
# SECTION 5: CRITICAL INSIGHT (Slide 16)
# ============================================================================

slide = add_content_slide(prs, "CRITICAL INSIGHT: Timing Affects Learning Speed, Not Valuation Level", COLOR_DARK_BLUE)
add_text_box(slide, 0.5, 1.2, 9, 0.35, "The Apparent Contradiction Resolved:", size=12, bold=True)

add_text_box(slide, 0.4, 1.7, 2.95, 0.3, "ESSAY 1: Returns", size=10, bold=True, color=COLOR_TEAL)
add_text_box(slide, 0.4, 2.05, 2.95, 2.3,
    "NO EFFECT\n\nTiming doesn't change what markets conclude\n\nValuation LEVEL unchanged",
    size=9)

add_text_box(slide, 3.5, 1.7, 2.95, 0.3, "ESSAY 2: Volatility", size=10, bold=True, color=COLOR_ORANGE)
add_text_box(slide, 3.5, 2.05, 2.95, 2.3,
    "INCREASES +1.83%**\n\nRegulatory constraint → forced incomplete disclosure\n\nLearning SPEED disrupted",
    size=9)

add_text_box(slide, 6.55, 1.7, 2.95, 0.3, "ESSAY 3: Turnover", size=10, bold=True, color=COLOR_RED)
add_text_box(slide, 6.55, 2.05, 2.95, 2.3,
    "ACCELERATES +5.3pp**\n\nStakeholder activation through immediate disclosure\n\nOrgizational RESPONSE faster",
    size=9)

add_text_box(slide, 0.5, 4.6, 9, 2.7,
    "THE KEY DISTINCTION: Timing affects LEARNING SPEED and STAKEHOLDER RESPONSE, not VALUATION LEVEL\n\nMarkets learn equally fast or slow (doesn't change CAR) • BUT learning process matters (uncertainty) • AND forced disclosure activates stakeholders (governance)\n\nAll three operating SIMULTANEOUSLY through DIFFERENT channels.",
    size=11, bold=True, color=RGBColor(192, 57, 43))

# ============================================================================
# SECTION 6: ESSAY 3 (Slides 17-19)
# ============================================================================

slide = add_content_slide(prs, "Essay 3: Governance Response - Executive Turnover", COLOR_RED)
bullets = [
    "Baseline executive turnover (all breaches, 30d): 46.4%",
    "With immediate disclosure: 50.6%  |  With delayed disclosure: 45.3%",
    "DIFFERENCE: +5.3 percentage points (p=0.008**)",
    " ",
    "Effect dynamics: PEAK at 30 days | DECAYS by 90 days | STABILIZES at 180 days",
    "→ Acute crisis response, not embedded governance reform",
    " ",
    "Sample: 896 breaches | Scope: Officers & directors | Mean: 3.2 executives per breach"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=12)

slide = add_content_slide(prs, "Essay 3: Stakeholder Activation Mechanism", COLOR_RED)
bullets = [
    "Step 1: Public Disclosure → Stakeholder Activation (investors, employees, regulators aware)",
    "Step 2: Board Response → Perceives breach as organizational crisis requiring visible action",
    "Step 3: Executive Turnover → Changes signal accountability, competence, fresh oversight",
    " ",
    "CRITICAL VALIDATION: Volatility does NOT mediate timing → turnover",
    "  • Mediation analysis (Script 91): Indirect effect = -0.0114 (NS)",
    "  • 95% CI includes zero → Governance response is DIRECT stakeholder pressure",
    "  • NOT driven by information quality effects (Essay 2)",
    "  • PROVES: Essays 2 and 3 are independent mechanisms, not cascade"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=11)

slide = add_content_slide(prs, "Essay 3: Size Heterogeneity", COLOR_RED)
bullets = [
    "Timing effect on turnover by firm size:",
    "  Q1 (Smallest): -0.679 (p=0.081, marginal)  |  Q2: -1.132** (p=0.026)",
    "  Q3 (Med-large): -1.651*** (p=0.006)  |  Q4 (Largest): +0.371 (NS)",
    " ",
    "Pattern: Medium firms show STRONGEST governance response | Largest firms unresponsive",
    " ",
    "NOTE: Size pattern DIFFERENT from Essay 2",
    "  Essay 2: Small firms hurt (+7.31%), large firms helped (-3.39%)",
    "  Essay 3: Medium firms respond most (logit effects)",
    "  → VALIDATES: Separate mechanisms, not spillover effects"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=11)

# ============================================================================
# SECTION 7: SYNTHESIS (Slides 20-22)
# ============================================================================

slide = add_content_slide(prs, "Cross-Essay Integration: Three Independent Mechanisms", COLOR_DARK_BLUE)
add_text_box(slide, 0.5, 1.2, 9, 0.35, "How Disclosure Requirements Work Through Multiple Pathways:", size=11, bold=True)

add_text_box(slide, 0.4, 1.7, 1.65, 0.3, "Essay 1", size=10, bold=True, color=COLOR_TEAL)
add_text_box(slide, 0.4, 2.05, 1.65, 2.1,
    "Stock Returns\n\nNO effect\n\nTiming irrelevant\n\nValuation Level (UNCHANGED)",
    size=9)

add_text_box(slide, 2.2, 1.7, 1.65, 0.3, "Essay 2", size=10, bold=True, color=COLOR_ORANGE)
add_text_box(slide, 2.2, 2.05, 1.65, 2.1,
    "Volatility\n\n+1.83%**\n\nIncomplete disclosure\n\nInfo Quality (DECREASED)",
    size=9)

add_text_box(slide, 4.0, 1.7, 1.65, 0.3, "Essay 3", size=10, bold=True, color=COLOR_RED)
add_text_box(slide, 4.0, 2.05, 1.65, 2.1,
    "Turnover\n\n+5.3pp**\n\nStakeholder pressure\n\nOrg. Response (ACCELERATED)",
    size=9)

add_text_box(slide, 5.8, 1.7, 3.8, 0.3, "Supporting Theory", size=10, bold=True, color=COLOR_GRAY)
add_text_box(slide, 5.8, 2.05, 3.8, 2.1,
    "Myers & Majluf (1984)\nKothari et al. (2009)\n\nTushman & Nadler (1978)\nMitchell, Agle, & Wood (1997)",
    size=9)

add_text_box(slide, 0.5, 4.3, 9, 2.9,
    "CENTRAL FINDING: Same regulation (FCC Rule 37.3) produces three distinct outcomes through three separate mechanisms. All operating SIMULTANEOUSLY.\n\nThis resolves apparent contradictions and explains policy complexity: regulation achieves some goals (governance) while creating unintended consequences (information quality loss) and leaving others unaffected (valuations).",
    size=11, bold=True, color=RGBColor(192, 57, 43))

slide = add_content_slide(prs, "Policy Implications: Evidence-Based & Scoped", COLOR_DARK_BLUE)
bullets = [
    "CLAIM 1: Stock Market Discipline Does NOT Operate Through Timing",
    "  Evidence: H1 null, TOST confirmed  |  SCOPE: Shareholder reactions only (not consumer protection, trust)",
    " ",
    "CLAIM 2: Timing Requirements Create Information Quality Tradeoffs",
    "  Evidence: +1.83%** volatility  |  SCOPE: FCC firms; effect strongest in small firms",
    " ",
    "CLAIM 3: Disclosure Successfully Activates Governance Response",
    "  Evidence: +5.3pp** turnover, not mediated by volatility  |  SCOPE: Executive/board changes; medium firms respond most",
    " ",
    "POLICY: Balance competing goals—disclosure speed (activates governance) vs. completeness (affects uncertainty)"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=10)

slide = add_content_slide(prs, "Limitations & Future Research", COLOR_DARK_BLUE)
bullets = [
    "Limitations: Stock market only | Publicly-traded firms | 2006-2025 | 30/90/180-day windows",
    " ",
    "Future: Cross-sector testing (SEC cybersecurity rule, HIPAA) | Consumer outcomes | Investigation quality",
    " ",
    "International comparison (GDPR) | Long-term governance reform | Direct measurement of investigation completeness",
    " ",
    "CLOSING: Disclosure policy is complex. Single-dimension theories fail.",
    "Move beyond 'faster is better' and design regulation balancing all three mechanisms."
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=11)

# ============================================================================
# APPENDIX (Slides 23-25)
# ============================================================================

slide = add_content_slide(prs, "Appendix: Robustness Testing & Data Quality", COLOR_GRAY)
bullets = [
    "H1 Timing tested across 27+ specifications (event windows, timing thresholds, subsamples, SE methods)",
    "ML validation: Random Forest & Gradient Boosting (timing not predictive)",
    " ",
    "Data Quality:",
    "  CRSP match: 87.9% of final sample (926 of 1,054) | 92.1% of raw extraction pool (~1,006)",
    "  CVSS scores: 80.4%  |  SEC 8-K data: 90.5%  |  Complete-case analysis (no imputation)",
    " ",
    "Economic Impact:",
    "  Empirical aggregate: $0.76B  |  Theoretical projection: $9.9B annual"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=11)

slide = add_content_slide(prs, "Appendix: Methodological References", COLOR_GRAY)
bullets = [
    "EVENT STUDY: Brown & Warner (1985) | Fama & French (1993) three-factor model",
    "REGRESSION: OLS HC3 standard errors | Logit with cluster SEs | Mediation with delta-method SEs",
    "IDENTIFICATION: Angrist & Pischke (2009) natural experiments | Parallel trends | Balance testing",
    " ",
    "THEORY:",
    "  Akerlof (1970): Adverse Selection  |  Spence (1973): Costly Signaling",
    "  Myers & Majluf (1984): Information Signaling  |  Diamond & Verrecchia (1991): Forced Disclosure Paradox",
    "  Tushman & Nadler (1978): Information Processing Capacity",
    "  Freeman (1984) / Mitchell, Agle, & Wood (1997): Stakeholder Theory & Salience"
]
add_bullet_points(slide, 0.5, 1.2, 9, 6, bullets, size=10)

# ============================================================================
# Save
# ============================================================================

output_path = 'Dissertation_Defense_Presentation_FINAL.pptx'
prs.save(output_path)

print("=" * 80)
print("SUCCESS: Final Dissertation Defense Presentation Created")
print("=" * 80)
print(f"\nFile: {output_path}")
print(f"Total Slides: {len(prs.slides)}")
print("\nUPDATES MADE:")
print("  [OK] Theoretical chain slide (Akerlof > Spence > Myers-Majluf > Dissertation)")
print("  [OK] Essay 1 heterogeneity slide (CVSS +6.27%, Media +7.08%)")
print("  [OK] Mitchell, Agle, & Wood (1997) citations for stakeholder theory")
print("  [OK] Kothari et al. (2009) bad news accumulation referenced")
print("  [OK] FCC mandate corrected to 7 days (not 30)")
print("  [OK] Data quality: CRSP 1,054 denominator (not 1,006)")
print("  [OK] Prior Breaches emphasized as STRONGEST effect")
print("  [OK] Mediation analysis interpretation added (independence proof)")
print("  [OK] Clean presentation - no clutter")
print("  [OK] 23-25 slides optimized for 35-40 minute defense")
print("\nREADY FOR APRIL 23, 2026 DEFENSE")
print("=" * 80)
